import os
import subprocess
import platform
import webbrowser
import curses
import sys
import runpy
import shutil
import platform
import ctypes
import tkinter as tk
from tkinter import messagebox
import json
import importlib.util
import glob
import platform
import socket
import psutil
import math
import time
import threading
import zipfile
import tarfile
import datetime
import importlib.util
import hashlib
import inspect
import uuid

_active_loops = {}


# at module top
dock_window = None
dock_thread = None

bg_window = None
bg_thread = None

_aliases = {}

# Save alias file in the same directory as this script
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ALIAS_FILE = os.path.join(BASE_DIR, "pynix_aliases.json")

def save_aliases():
    """Save aliases to a JSON file (same directory as terminal)."""
    try:
        with open(ALIAS_FILE, "w") as f:
            json.dump(_aliases, f, indent=2)
    except Exception as e:
        print(f"⚠️ Failed to save aliases: {e}")

def load_aliases():
    """Load aliases from disk and register them as commands."""
    global _aliases
    if not os.path.exists(ALIAS_FILE):
        return
    try:
        with open(ALIAS_FILE, "r") as f:
            _aliases.update(json.load(f))
        for name, cmd in _aliases.items():
            registered_commands[name] = lambda a, c=cmd: handle_command(c)
        if _aliases:
            print(f"📦 Loaded {len(_aliases)} aliases from {ALIAS_FILE}")
    except Exception as e:
        print(f"⚠️ Failed to load aliases: {e}")
# Global state flag for the spinning donut
try:
    import msvcrt  # Windows key input
except ImportError:
    import termios, tty, select  # Unix key input

DonutRunning = False

# ==========================
# GLOBALS
# ==========================
command_history = []


# Global clipboard variable
clipboard = {"path": None, "is_folder": False}

"""Display a confirmation popup using Tkinter."""
def confirm_action(title, message):
    root = tk.Tk()
    root.withdraw()
    result = messagebox.askyesno(title, message)
    root.destroy()
    return result


    
# extneral command loader    
import traceback

def load_external_commands():
    """Load all Python command files from /commands and /commands/added."""
    base_dir = os.path.join(os.getcwd(), "commands")
    added_dir = os.path.join(base_dir, "added")

    os.makedirs(base_dir, exist_ok=True)
    os.makedirs(added_dir, exist_ok=True)

    loaded_count = 0
    failed_files = []

    # Scan both main /commands and /commands/added
    command_dirs = [base_dir, added_dir]

    for cmd_dir in command_dirs:
        if not os.path.exists(cmd_dir):
            continue

        for file in os.listdir(cmd_dir):
            if not file.endswith(".py"):
                continue

            file_path = os.path.join(cmd_dir, file)
            module_name = f"cmd_{os.path.basename(cmd_dir)}_{file[:-3]}"

            try:
                # Unload if already imported
                if module_name in sys.modules:
                    del sys.modules[module_name]

                spec = importlib.util.spec_from_file_location(module_name, file_path)
                if not spec or not spec.loader:
                    print(f"⚠️ Skipping invalid module: {file}")
                    continue

                # Create the module
                module = importlib.util.module_from_spec(spec)

                # Inject shared globals so register_command() works
                import __main__
                module.__dict__["registered_commands"] = __main__.registered_commands
                module.__dict__["handle_command"] = __main__.handle_command

                # Execute the file in that context
                spec.loader.exec_module(module)
                sys.modules[module_name] = module
                loaded_count += 1

            except Exception as e:
                failed_files.append((file, str(e)))
                traceback.print_exc()

    print(f"🔄 Loaded {loaded_count} external command(s) from /commands.")
    if failed_files:
        print("⚠️ Failed to load:")
        for name, err in failed_files:
            print(f"  - {name}: {err}")

    import glob

    def load_from_folder(folder):
        for file_path in glob.glob(os.path.join(folder, "*.py")):
            try:
                module_name = os.path.splitext(os.path.basename(file_path))[0]
                spec = importlib.util.spec_from_file_location(module_name, file_path)
                module = importlib.util.module_from_spec(spec)
                spec.loader.exec_module(module)
                print(f"✅ Loaded command: {module_name}")
            except Exception as e:
                print(f"⚠️ Failed to load {file_path}: {e}")

    load_from_folder(base_dir)
    load_from_folder(added_dir)

# Store last program list for launch by number
last_program_list = []

# =======================================
# Command registry
# =======================================
registered_commands = {}

def register_command(name):
    def wrapper(func):
        registered_commands[name] = func
        return func
    return wrapper


# =======================================
# Custom Commands
# =======================================

@register_command("whoami")
def whoami(args):
    """Display the current logged-in username (like Unix 'whoami')."""
    import getpass
    import platform
    import os

    try:
        user = getpass.getuser()
    except Exception:
        # Fallback for some edge environments
        user = os.environ.get("USERNAME") or os.environ.get("USER") or "Unknown"

    system = platform.system()
    print(f"{user} ({system})")


@register_command("hello")
def hello(args):
    print("Hello, world!")

@register_command("ls")
def list_files(args):
    path = args[0] if args else "."
    for f in os.listdir(path):
        print(f)

@register_command("clear")
def clear(args):
    os.system("cls" if os.name == "nt" else "clear")

last_dir = None

@register_command("goto")
def goto_cmd(args):
    import os
    global last_dir

    if not args:
        print("Usage: goto <path>")
        return

    target = os.path.abspath(os.path.expanduser(args[0]))

    if not os.path.exists(target):
        print(f"❌ Path not found: {target}")
        return

    if os.path.isdir(target):
        try:
            last_dir = os.getcwd()
            os.chdir(target)
            print(f"📁 Changed directory to: {target}")
        except Exception as e:
            print(f"⚠️ Error: {e}")
    else:
        print(f"📄 Target is a file: {target}")


@register_command("goback")
def goback(args):
    """Go back to the last visited directory (like 'cd -' in Unix)."""
    import os
    global last_dir

    if not last_dir:
        print("⚠️ No previous directory recorded yet.")
        return

    try:
        current = os.getcwd()
        os.chdir(last_dir)
        print(f"🔙 Moved back to: {os.getcwd()}")
        last_dir = current  # swap so you can toggle back and forth
    except Exception as e:
        print(f"⚠️ Error returning to last directory: {e}")



@register_command("cd")
def cd_cmd(args):
    """Change directory, remembering previous path for 'goback'."""
    import os
    global last_dir

    if not args:
        print(os.getcwd())
        return

    path = os.path.expanduser(args[0])
    path = os.path.abspath(path)

    if not os.path.exists(path):
        print(f"❌ Path not found: {path}")
        return

    if os.path.isdir(path):
        try:
            last_dir = os.getcwd()
            os.chdir(path)
            print(f"📁 Changed directory to: {path}")
        except Exception as e:
            print(f"⚠️ Error: {e}")
    else:
        print(f"📄 Target is a file: {path}")

@register_command("opendrop")
def opendrop(args):
    """Run the OpenDrop Python program."""
    if os.name != "nt":
        print("opendrop only works on Windows.")
        return

    script = "opendropUnix.py"
    if not os.path.exists(script):
        print(f"Error: '{script}' not found in the current directory.")
        return

    try:
        subprocess.call(["cmd", "/c", f"python {script}"])
    except Exception as e:
        print(f"Error running {script}: {e}")
        
@register_command("findgit")
def opendrop(args):

    if os.name != "nt":
        print("Finding Git Repos.")
        return

    script = "gitfind.py"
    if not os.path.exists(script):
        print(f"Error: '{script}' not found in the current directory.")
        return

    try:
        subprocess.call(["cmd", "/c", f"python {script}"])
    except Exception as e:
        print(f"Error running {script}: {e}")
        
@register_command("programs")
def list_programs(args):
    """List all .py and .html files in the current directory with numbers."""
    global last_program_list
    files = [f for f in os.listdir(".") if f.endswith(".py") or f.endswith(".html")]
    if not files:
        print("No .py or .html files found in this directory.")
        return

    last_program_list = files
    print("\nAvailable programs:")
    for i, f in enumerate(files, start=1):
        print(f"{i}. {f}")
    print("\nUse 'launch <number>' or 'launch <filename>' to open.\n")


@register_command("launch")
def launch_program(args):
    """
    Launch a program by number or filename.

    Usage:
        launch <number|filename>
        launch -n <number|filename>   # open in new terminal window

    Description:
        Launches a file or script from the last listed program list
        or by direct path. The -n flag opens it in a new PyShell terminal window.
    """
    import os
    import subprocess
    import webbrowser
    import sys

    global last_program_list

    if not args:
        print("Usage: launch <number or filename> [-n for new terminal window]")
        return

    # --- Detect -n flag ---
    new_window = False
    if "-n" in args:
        new_window = True
        args.remove("-n")

    if not args:
        print("Usage: launch <number or filename>")
        return

    target = args[0]

    # --- Resolve file target ---
    if target.isdigit():
        idx = int(target) - 1
        if 0 <= idx < len(last_program_list):
            file_to_open = last_program_list[idx]
        else:
            print("❌ Invalid number.")
            return
    else:
        file_to_open = target

    if not os.path.exists(file_to_open):
        print(f"❌ File not found: {file_to_open}")
        return

    # --- Launch behavior ---
    try:
        if file_to_open.endswith(".py"):
            if new_window:
                # 🪟 Launch in a new terminal window (cross-platform)
                if os.name == "nt":
                    subprocess.Popen(
                        ["start", "cmd", "/k", f'python "{file_to_open}"'],
                        shell=True
                    )
                else:
                    subprocess.Popen(
                        ["x-terminal-emulator", "-e", f'python3 "{file_to_open}"'],
                        shell=False
                    )
                print(f"🚀 Launched {file_to_open} in new terminal window.")
            else:
                subprocess.Popen([sys.executable, file_to_open], shell=True)
                print(f"🚀 Launched {file_to_open} (background).")

        elif file_to_open.endswith(".html"):
            webbrowser.open(file_to_open)
            print(f"🌐 Opening {file_to_open} in web browser.")

        else:
            os.startfile(file_to_open)
            print(f"📂 Opening {file_to_open} with default program.")

    except Exception as e:
        print(f"⚠️ Error launching file: {e}")

        
@register_command("touch")
def touch(args):
    """Create a new file with the specified name and extension."""
    if not args:
        print("Usage: touch <filename>.<extension>")
        return

    filename = args[0]
    full_path = os.path.join(os.getcwd(), filename)

    try:
        # Create file if it doesn't exist, or update timestamp if it does
        with open(full_path, 'a'):
            os.utime(full_path, None)
        print(f"Created: {full_path}")
    except Exception as e:
        print(f"Error creating file: {e}")
        
@register_command("rm")
def remove_item(args):
    """Remove one or more files or folders."""
    if not args:
        print("Usage: rm <name> [name2 name3 ...]")
        return

    for name in args:
        path = os.path.join(os.getcwd(), name)

        if not os.path.exists(path):
            print(f"Not found: {name}")
            continue

        try:
            if os.path.isfile(path):
                os.remove(path)
                print(f"Removed file: {name}")
            elif os.path.isdir(path):
                # Delete folders (including non-empty)
                shutil.rmtree(path)
                print(f"Removed folder: {name}")
            else:
                print(f"Unknown type (skipped): {name}")
        except Exception as e:
            print(f"Error removing {name}: {e}")

import curses

@register_command("nano")
def nano(args):
    """Advanced Nano-like editor with multiple files, clipboard, and top/bottom navigation."""
    import curses, pyperclip, os

    start_file = args[0] if args else None
    open_files = []

    def load_file(path):
        lines = []
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                lines = f.read().splitlines()
        else:
            lines = [""]
        return {"name": path, "lines": lines, "y": 0, "x": 0, "scroll": 0}

    # Start with at least one file
    if start_file:
        open_files.append(load_file(os.path.abspath(start_file)))
    else:
        open_files.append({"name": "untitled.txt", "lines": [""], "y": 0, "x": 0, "scroll": 0})

    current = 0  # active file index

    def editor(stdscr):
        nonlocal current
        curses.curs_set(1)
        stdscr.keypad(True)

        while True:
            file = open_files[current]
            lines, y, x, scroll = file["lines"], file["y"], file["x"], file["scroll"]
            stdscr.clear()
            h, w = stdscr.getmaxyx()
            num_width = 6
            content_width = w - num_width

            # === Top bar: all open files ===
            bar = " | ".join(
                [f"[{os.path.basename(f['name'])}]" if i == current else os.path.basename(f["name"])
                 for i, f in enumerate(open_files)]
            )
            stdscr.addstr(0, 0, bar[:w - 1], curses.A_REVERSE)

            # === Controls header ===
            header = "[Ctrl+S Save] [Ctrl+Q Quit] [Ctrl+W Close tab] [Ctrl+V Paste] [Ctrl+T New File] [Ctrl+O Open] [Tab Switch] [Ctrl+G Top] [Ctrl+J Bottom]"
            stdscr.addstr(1, 0, header[:w - 1], curses.A_REVERSE)

            # === Draw file contents ===
            visible = lines[scroll:scroll + h - 3]
            for i, line in enumerate(visible):
                ln = scroll + i + 1
                ln_str = f"{ln:4d} | "
                try:
                    stdscr.addstr(i + 2, 0, ln_str)
                    stdscr.addstr(i + 2, num_width, line[:content_width])
                except curses.error:
                    pass

            draw_y = y - scroll + 2
            if 2 <= draw_y < h - 1:
                stdscr.move(draw_y, min(x + num_width, w - 2))
            stdscr.refresh()

            key = stdscr.getch()

            # Quit
            if key == 17:  # Ctrl+Q
                break

            # Save
            elif key == 19:  # Ctrl+S
                with open(file["name"], "w", encoding="utf-8") as f:
                    f.write("\n".join(lines))
                stdscr.addstr(h - 1, 0, f"💾 Saved {file['name']}".ljust(w - 1), curses.A_REVERSE)
                stdscr.refresh(); curses.napms(400)

            # Paste
            elif key == 22:  # Ctrl+V
                try:
                    clip = pyperclip.paste()
                    if clip:
                        for line in clip.splitlines():
                            lines[y] = lines[y][:x] + line + lines[y][x:]
                            y += 1
                            lines.insert(y, "")
                        stdscr.addstr(h - 1, 0, "📋 Pasted".ljust(w - 1), curses.A_REVERSE)
                        stdscr.refresh(); curses.napms(300)
                except Exception as e:
                    stdscr.addstr(h - 1, 0, f"⚠️ Paste failed: {e}".ljust(w - 1), curses.A_REVERSE)
                    stdscr.refresh(); curses.napms(600)

            # Open new file
            elif key == 15:  # Ctrl+O
                curses.echo()
                stdscr.addstr(h - 1, 0, "Open file: ".ljust(w - 1), curses.A_REVERSE)
                stdscr.clrtoeol()
                path = stdscr.getstr(h - 1, 12, 100).decode().strip()
                curses.noecho()
                if path:
                    try:
                        newf = load_file(os.path.abspath(path))
                        open_files.append(newf)
                        current = len(open_files) - 1
                        stdscr.addstr(h - 1, 0, f"📂 Opened {path}".ljust(w - 1), curses.A_REVERSE)
                        stdscr.refresh(); curses.napms(400)
                    except Exception as e:
                        stdscr.addstr(h - 1, 0, f"❌ Error: {e}".ljust(w - 1), curses.A_REVERSE)
                        stdscr.refresh(); curses.napms(800)
             
            # Create new file
            elif key == 20:  # Ctrl+T
                curses.echo()
                stdscr.addstr(h - 1, 0, "New file name: ".ljust(w - 1), curses.A_REVERSE)
                stdscr.clrtoeol()
                path = stdscr.getstr(h - 1, 15, 100).decode().strip()
                curses.noecho()
                if path:
                    try:
                        abs_path = os.path.abspath(path)
                        # Create the empty file if it doesn’t exist
                        if not os.path.exists(abs_path):
                            with open(abs_path, "w", encoding="utf-8") as f:
                                f.write("")
                        newf = load_file(abs_path)
                        open_files.append(newf)
                        current = len(open_files) - 1
                        stdscr.addstr(h - 1, 0, f"📝 Created {path}".ljust(w - 1), curses.A_REVERSE)
                        stdscr.refresh(); curses.napms(400)
                    except Exception as e:
                        stdscr.addstr(h - 1, 0, f"❌ Error: {e}".ljust(w - 1), curses.A_REVERSE)
                        stdscr.refresh(); curses.napms(800)
             
                         # Close current file
            elif key == 23:  # Ctrl+W
                if len(open_files) > 1:
                    closed = open_files.pop(current)
                    current = max(0, current - 1)
                    stdscr.addstr(h - 1, 0,
                                  f"❌ Closed {os.path.basename(closed['name'])}".ljust(w - 1),
                                  curses.A_REVERSE)
                    stdscr.refresh(); curses.napms(400)
                else:
                    stdscr.addstr(h - 1, 0,
                                  "⚠️ Cannot close last open file".ljust(w - 1),
                                  curses.A_REVERSE)
                    stdscr.refresh(); curses.napms(600)
           

            # Switch between files
            elif key == 9:  # Tab
                if len(open_files) > 1:
                    current = (current + 1) % len(open_files)
                    stdscr.addstr(h - 1, 0,
                                  f"➡️  Switched to {open_files[current]['name']}".ljust(w - 1),
                                  curses.A_REVERSE)
                    stdscr.refresh(); curses.napms(300)

            # Jump to Top
            elif key == 7:  # Ctrl+H
                y, x, scroll = 0, 0, 0
                stdscr.addstr(h - 1, 0, "⬆️  Jumped to top".ljust(w - 1), curses.A_REVERSE)
                stdscr.refresh(); curses.napms(300)

            # Jump to Bottom
            elif key == 10:  # Ctrl+J
                y = len(lines) - 1
                x = len(lines[y])
                scroll = max(0, len(lines) - (h - 3))
                stdscr.addstr(h - 1, 0, "⬇️  Jumped to bottom".ljust(w - 1), curses.A_REVERSE)
                stdscr.refresh(); curses.napms(300)

            # Navigation
            elif key == curses.KEY_UP:
                if y > 0: y -= 1
                elif scroll > 0: scroll -= 1
            elif key == curses.KEY_DOWN:
                if y < len(lines) - 1: y += 1
                else: lines.append("")
                if y - scroll >= h - 4: scroll += 1
            elif key == curses.KEY_LEFT:
                if x > 0: x -= 1
                elif y > 0: y -= 1; x = len(lines[y])
            elif key == curses.KEY_RIGHT:
                if x < len(lines[y]): x += 1
                elif y < len(lines) - 1: y += 1; x = 0

            # Backspace
            elif key in (8, 127, curses.KEY_BACKSPACE):
                if x > 0:
                    lines[y] = lines[y][:x - 1] + lines[y][x:]
                    x -= 1
                elif y > 0:
                    prev_len = len(lines[y - 1])
                    lines[y - 1] += lines[y]
                    del lines[y]
                    y -= 1; x = prev_len

            # Enter
            elif key in (13,):
                new_line = lines[y][x:]
                lines[y] = lines[y][:x]
                lines.insert(y + 1, new_line)
                y += 1; x = 0

            # Printable text
            elif 32 <= key <= 126:
                lines[y] = lines[y][:x] + chr(key) + lines[y][x:]
                x += 1

            # Save state
            file.update({"lines": lines, "y": y, "x": x, "scroll": scroll})

    try:
        curses.wrapper(editor)
    except Exception as e:
        print(f"Editor error: {e}")
        
@register_command("rep")
def rep_cmd(args):
    """
    Replace one file with another.
    
    Usage:
      rep <source> -t <target>

    Example:
      rep new_version.py -t main.py
    """
    import os
    import shutil

    if len(args) < 3 or args[1] != "-t":
        print("Usage: rep <source> -t <target>")
        return

    source, target = args[0], args[2]
    source = os.path.abspath(source)
    target = os.path.abspath(target)

    if not os.path.exists(source):
        print(f"❌ Source file not found: {source}")
        return

    # Create target directory if it doesn't exist
    target_dir = os.path.dirname(target)
    if target_dir and not os.path.exists(target_dir):
        try:
            os.makedirs(target_dir)
            print(f"📁 Created missing target directory: {target_dir}")
        except Exception as e:
            print(f"⚠️ Failed to create target directory: {e}")
            return

    try:
        shutil.copy2(source, target)
        print(f"✅ Replaced '{target}' with '{source}' successfully.")
    except PermissionError:
        print("⚠️ Permission denied while replacing file.")
    except Exception as e:
        print(f"⚠️ Replacement failed: {e}")

@register_command("gitpush")
def gitpush(args):
    """
    Runs 'git push -u origin main' in the system shell.
    
    Usage:
      gitpush
    
    Description:
      Executes a real Git push to the current repository’s remote origin main branch.
      Works on Windows, macOS, and Linux.
    """
    import os
    import subprocess
    import platform

    print("🚀 Running: git push -u origin main\n")

    try:
        if platform.system() == "Windows":
            # Windows uses PowerShell or cmd
            subprocess.run(["git", "push", "-u", "origin", "main"], shell=True)
        else:
            # macOS / Linux
            subprocess.run("git push -u origin main", shell=True, executable="/bin/bash")
        print("\n✅ Push complete.")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Push failed: {e}")

@register_command("gitpull")
def gitpull_cmd(args):
    """
    Downloads a GitHub repository or a specific file from a public repo.

    Usage:
      gitpull <username>/<repository>
      gitpull <username>/<repository>/<file>

    Examples:
      gitpull torvalds/linux
      gitpull torvalds/linux/README.md

    Description:
      Pulls files or entire repos directly from GitHub without requiring Git.
      Automatically saves files into the current directory.
    """
    import os
    import urllib.request
    import zipfile
    import io

    if not args:
        print("Usage: gitpull <username>/<repo> or gitpull <username>/<repo>/<file>")
        return

    target = args[0].strip("/")
    parts = target.split("/")

    if len(parts) < 2:
        print("❌ Invalid format. Use: gitpull <username>/<repo> or gitpull <username>/<repo>/<file>")
        return

    username, repo = parts[0], parts[1]
    file_path = "/".join(parts[2:]) if len(parts) > 2 else None

    if file_path:
        # --- Pull specific file ---
        url = f"https://raw.githubusercontent.com/{username}/{repo}/main/{file_path}"
        save_path = os.path.join(os.getcwd(), os.path.basename(file_path))
        print(f"⬇️ Downloading file: {url}")
        try:
            urllib.request.urlretrieve(url, save_path)
            print(f"✅ Saved file: {save_path}")
        except Exception as e:
            print(f"⚠️ Failed to download file: {e}")
    else:
        # --- Pull entire repository ---
        url = f"https://github.com/{username}/{repo}/archive/refs/heads/main.zip"
        print(f"⬇️ Downloading repository: {url}")
        try:
            with urllib.request.urlopen(url) as response:
                data = response.read()
            with zipfile.ZipFile(io.BytesIO(data)) as zip_ref:
                zip_ref.extractall(os.getcwd())
            print(f"✅ Repository '{repo}' extracted successfully into current directory.")
        except Exception as e:
            print(f"⚠️ Failed to download repository: {e}")

@register_command("gitcommit")
def gitcommit(args):
    """
    Runs 'git commit -m "auto commit"' in the system shell.

    Usage:
      gitcommit

    Description:
      Automatically commits all staged changes with the message 'auto commit'.
      Displays the full Git output in PyNixShell.
    """
    import subprocess
    import platform

    print("💾 Running: git commit -m \"auto commit\"\n")

    try:
        # Cross-platform execution
        if platform.system() == "Windows":
            subprocess.run(["git", "commit", "-m", "auto commit"], shell=True)
        else:
            subprocess.run("git commit -m 'auto commit'", shell=True, executable="/bin/bash")

        print("\n✅ Commit complete.")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Commit failed: {e}")

@register_command("gitstage")
def gitstage(args):
    """
    Stages all changes or a specific file using Git.

    Usage:
      gitstage            → stages all changes (git add .)
      gitstage <filename> → stages a specific file

    Example:
      gitstage
      gitstage main.py
    """
    import subprocess
    import platform

    if args:
        target = args[0]
        cmd = ["git", "add", target]
        print(f"📁 Running: git add {target}\n")
    else:
        cmd = ["git", "add", "."]
        print("📦 Running: git add . (stage all changes)\n")

    try:
        if platform.system() == "Windows":
            subprocess.run(cmd, shell=True)
        else:
            subprocess.run(" ".join(cmd), shell=True, executable="/bin/bash")

        print("✅ Files staged successfully.")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Stage failed: {e}")


        
@register_command("rfpt")
def refresh_pyterm(args):
    """Re-executes PyTerm code in place, updating all commands without closing."""
    print("Refreshing PyTerm in place...")

    script_path = os.path.abspath(sys.argv[0])
    try:
        # Re-run the script file (like restarting, but inside same interpreter)
        runpy.run_path(script_path, run_name="__main__")
        print("PyTerm refreshed successfully.")
        handle_command("clear")
    except Exception as e:
        print(f"Error refreshing PyTerm: {e}")   


   
   
        
@register_command("gitrm")
def gitrm(args):
    """
    Removes a file from the Git repository index but keeps it locally.

    Usage:
      gitrm <filename>

    Example:
      gitrm secrets.txt

    Description:
      Untracks a file from Git while keeping it in your folder.
      Equivalent to: git rm --cached <file>
    """
    import subprocess
    import platform

    if not args:
        print("Usage: gitrm <filename>")
        return

    filename = args[0]
    print(f"🗑️ Untracking file: {filename}\n")

    try:
        if platform.system() == "Windows":
            subprocess.run(["git", "rm", "--cached", filename], shell=True)
        else:
            subprocess.run(f"git rm --cached '{filename}'", shell=True, executable="/bin/bash")

        print(f"✅ '{filename}' removed from repository index (kept locally).")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Removal failed: {e}")
        
        
@register_command("gitinit")
def gitinit(args):
    """
    Initializes a real Git repository using system Git.

    Usage:
      gitinit

    Description:
      Runs 'git init' in the current working directory to create a new Git repository.
      Displays Git’s real output inline in PyNixShell.
    """
    import subprocess
    import platform

    print("🧩 Running: git init\n")

    try:
        if platform.system() == "Windows":
            subprocess.run(["git", "init"], shell=True)
        else:
            subprocess.run("git init", shell=True, executable="/bin/bash")

        print("\n✅ Repository initialized successfully.")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Initialization failed: {e}")
        
@register_command("gitstatus")
def gitstatus(args):
    """
    Displays a clean, color-coded summary of your current Git repository status.

    Usage:
      gitstatus

    Description:
      Runs 'git status --porcelain' and formats the output for readability.
      Shows clear sections for modified, new, deleted, and untracked files.
    """
    import subprocess
    import platform

    # Simple color helper (ANSI codes)
    def color(text, code): return f"\033[{code}m{text}\033[0m"

    print("📋 Checking repository status...\n")

    try:
        # Run git status in porcelain mode
        result = subprocess.run(
            ["git", "status", "--porcelain"],
            capture_output=True,
            text=True,
            shell=(platform.system() == "Windows")
        )

        output = result.stdout.strip()
        if not output:
            print(color("✅ Working directory clean — nothing to commit.", "32"))
            return

        modified, added, deleted, untracked = [], [], [], []

        for line in output.splitlines():
            status = line[:2].strip()
            filename = line[3:].strip() if len(line) > 3 else line.strip()

            if status in ("M", "MM", "AM", "MA"):
                modified.append(filename)
            elif status in ("A", "??A"):
                added.append(filename)
            elif status in ("D", "AD", "MD"):
                deleted.append(filename)
            elif status == "??":
                untracked.append(filename)

        # Display sections neatly
        if added:
            print(color("🟩 Added files:", "32"))
            for f in added:
                print(f"  + {f}")
            print()

        if modified:
            print(color("🟨 Modified files:", "33"))
            for f in modified:
                print(f"  * {f}")
            print()

        if deleted:
            print(color("🟥 Deleted files:", "31"))
            for f in deleted:
                print(f"  - {f}")
            print()

        if untracked:
            print(color("⚪ Untracked files:", "37"))
            for f in untracked:
                print(f"  ? {f}")
            print()

        print(color("💡 Tip:", "36"), "Use 'gitstage' or 'gitadd' to stage files.")
    except FileNotFoundError:
        print("❌ Git is not installed or not found in PATH.")
    except Exception as e:
        print(f"⚠️ Error running gitstatus: {e}")
        
        
@register_command("cat")
def cat(args):
    """Display the contents of a text file."""
    if not args:
        print("Usage: cat <filename>")
        return

    filename = args[0]
    path = os.path.join(os.getcwd(), filename)

    if not os.path.exists(path):
        print(f"File not found: {filename}")
        return

    if os.path.isdir(path):
        print(f"'{filename}' is a directory.")
        return

    try:
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
            print(content)
    except Exception as e:
        print(f"Error reading file: {e}")
        
@register_command("back")
def back(args):
    """Go back one directory (equivalent to 'cd ..')."""
    try:
        os.chdir("..")
        print(f"Moved back to: {os.getcwd()}")
    except Exception as e:
        print(f"Error moving back: {e}")
        
@register_command("mkdir")
def make_dir(args):
    """Create a new directory (like Unix mkdir)."""
    if not args:
        print("Usage: mkdir <foldername> or mkdir -p <path>")
        return

    # Handle optional -p flag
    if args[0] == "-p":
        if len(args) < 2:
            print("Usage: mkdir -p <path>")
            return
        path = args[1]
        full_path = os.path.join(os.getcwd(), path)
        try:
            os.makedirs(full_path, exist_ok=True)
            print(f"Created directory (with parents): {full_path}")
        except Exception as e:
            print(f"Error creating directory: {e}")
    else:
        path = args[0]
        full_path = os.path.join(os.getcwd(), path)
        try:
            os.mkdir(full_path)
            print(f"Created directory: {full_path}")
        except FileExistsError:
            print(f"Directory already exists: {path}")
        except Exception as e:
            print(f"Error creating directory: {e}")
            
@register_command("tree")
def tree(args):
    """Display a directory tree (Unix-style)."""
    import shutil

    # Parse arguments
    path = os.getcwd()
    max_depth = None

    if args:
        if args[0] == "-L" and len(args) > 1 and args[1].isdigit():
            max_depth = int(args[1])
            if len(args) > 2:
                path = args[2]
        else:
            path = args[0]

    if not os.path.exists(path):
        print(f"Path not found: {path}")
        return

    # Determine terminal width for wrapping
    term_width = shutil.get_terminal_size((80, 20)).columns

    print(path)

    def walk(dir_path, prefix="", depth=0):
        if max_depth is not None and depth >= max_depth:
            return

        try:
            entries = sorted(os.listdir(dir_path))
        except PermissionError:
            print(prefix + "[Permission Denied]")
            return

        dirs = [e for e in entries if os.path.isdir(os.path.join(dir_path, e))]
        files = [e for e in entries if os.path.isfile(os.path.join(dir_path, e))]
        entries = dirs + files

        for i, entry in enumerate(entries):
            full_path = os.path.join(dir_path, entry)
            connector = "└── " if i == len(entries) - 1 else "├── "
            line = prefix + connector + entry
            print(line[:term_width])

            if os.path.isdir(full_path):
                extension = "    " if i == len(entries) - 1 else "│   "
                walk(full_path, prefix + extension, depth + 1)

    walk(path)
            
@register_command("du")
def du(args):
    """Show disk usage (like Unix 'du'). Supports -s for summary and -f for file."""
    import math

    path = os.getcwd()
    show_total_only = False
    show_file_only = False
    file_target = None

    # Parse arguments
    if args:
        if args[0] == "-s":
            show_total_only = True
            if len(args) > 1:
                path = args[1]
        elif args[0] == "-f" and len(args) > 1:
            show_file_only = True
            file_target = args[1]
        else:
            path = args[0]

    def human_readable(size_bytes):
        if size_bytes == 0:
            return "0 B"
        size_name = ("B", "KB", "MB", "GB", "TB")
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {size_name[i]}"

    def get_size(start_path):
        total_size = 0
        for dirpath, dirnames, filenames in os.walk(start_path):
            for f in filenames:
                fp = os.path.join(dirpath, f)
                if not os.path.islink(fp):
                    try:
                        total_size += os.path.getsize(fp)
                    except FileNotFoundError:
                        pass
        return total_size

    # --- Handle file mode ---
    if show_file_only:
        file_path = os.path.join(os.getcwd(), file_target)
        if not os.path.exists(file_path):
            print(f"File not found: {file_target}")
            return
        if os.path.isdir(file_path):
            print(f"'{file_target}' is a directory. Use du -s instead.")
            return
        size = os.path.getsize(file_path)
        print(f"{human_readable(size)}\t{file_target}")
        return

    # --- Handle folder modes ---
    if not os.path.exists(path):
        print(f"Path not found: {path}")
        return

    if show_total_only:
        total = get_size(path)
        print(f"{human_readable(total)}\t{path}")
    else:
        for root, dirs, files in os.walk(path):
            size = get_size(root)
            print(f"{human_readable(size)}\t{root}")
     
@register_command("df")
def disk_free(args):
    """Show disk space usage (like Unix 'df')."""
    import shutil
    import math

    def human_readable(size_bytes):
        """Convert bytes to human-readable format."""
        if size_bytes == 0:
            return "0 B"
        size_name = ("B", "KB", "MB", "GB", "TB")
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {size_name[i]}"

    human = "-h" in args

    print("Filesystem".ljust(20), "Size".rjust(10), "Used".rjust(10),
          "Avail".rjust(10), "Use%".rjust(8), "Mounted on")

    if os.name == "nt":
        # Windows: list all available drives
        import string
        from ctypes import windll

        bitmask = windll.kernel32.GetLogicalDrives()
        for letter in string.ascii_uppercase:
            if bitmask & 1:
                drive = f"{letter}:\\"
                try:
                    total, used, free = shutil.disk_usage(drive)
                    use_percent = (used / total) * 100 if total > 0 else 0
                    print(
                        drive.ljust(20),
                        (human_readable(total) if human else str(total)).rjust(10),
                        (human_readable(used) if human else str(used)).rjust(10),
                        (human_readable(free) if human else str(free)).rjust(10),
                        f"{use_percent:.0f}%".rjust(8),
                        drive
                    )
                except Exception:
                    pass
            bitmask >>= 1
    else:
        # Unix/macOS
        mounts = ["/"]
        for m in mounts:
            total, used, free = shutil.disk_usage(m)
            use_percent = (used / total) * 100 if total > 0 else 0
            print(
                m.ljust(20),
                (human_readable(total) if human else str(total)).rjust(10),
                (human_readable(used) if human else str(used)).rjust(10),
                (human_readable(free) if human else str(free)).rjust(10),
                f"{use_percent:.0f}%".rjust(8),
                m
            )
 
@register_command("cp")
def copy_item(args):
    """Copy a file or folder into memory for later pasting."""
    if not args:
        print("Usage: cp <filename or foldername>")
        return

    target = args[0]
    path = os.path.join(os.getcwd(), target)

    if not os.path.exists(path):
        print(f"Not found: {target}")
        return

    clipboard["path"] = path
    clipboard["is_folder"] = os.path.isdir(path)
    print(f"{'Folder' if clipboard['is_folder'] else 'File'} copied: {target}")


@register_command("paste")
def paste_item(args):
    """Paste the last copied file or folder into the current directory."""
    if not clipboard["path"]:
        print("Nothing copied. Use 'cp <filename>' first.")
        return

    src = clipboard["path"]
    dst_name = os.path.basename(src)
    dst_path = os.path.join(os.getcwd(), dst_name)

    # Prevent overwriting by adding (copy) if it already exists
    if os.path.exists(dst_path):
        base, ext = os.path.splitext(dst_name)
        dst_name = f"{base} (copy){ext}"
        dst_path = os.path.join(os.getcwd(), dst_name)

    try:
        if clipboard["is_folder"]:
            shutil.copytree(src, dst_path)
            print(f"Pasted folder: {dst_name}")
        else:
            shutil.copy2(src, dst_path)
            print(f"Pasted file: {dst_name}")
    except Exception as e:
        print(f"Error pasting item: {e}") 

@register_command("head")
def head(args):
    """Display the first few lines of a file (default 3, like Unix head)."""
    if not args:
        print("Usage: head [-n num] <filename>")
        return

    num_lines = 3  # default
    filename = None

    # Parse arguments
    if args[0] == "-n":
        if len(args) < 3 or not args[1].isdigit():
            print("Usage: head -n <number> <filename>")
            return
        num_lines = int(args[1])
        filename = args[2]
    else:
        filename = args[0]

    path = os.path.join(os.getcwd(), filename)

    if not os.path.exists(path):
        print(f"File not found: {filename}")
        return

    if os.path.isdir(path):
        print(f"'{filename}' is a directory.")
        return

    try:
        with open(path, "r", encoding="utf-8") as f:
            for i, line in enumerate(f):
                if i >= num_lines:
                    break
                print(line.rstrip())
    except Exception as e:
        print(f"Error reading file: {e}")
        
@register_command("tail")
def tail(args):
    """Display the last few lines of a file (default 3, like Unix tail)."""
    if not args:
        print("Usage: tail [-n num] <filename>")
        return

    num_lines = 3  # default
    filename = None

    # Parse arguments
    if args[0] == "-n":
        if len(args) < 3 or not args[1].isdigit():
            print("Usage: tail -n <number> <filename>")
            return
        num_lines = int(args[1])
        filename = args[2]
    else:
        filename = args[0]

    path = os.path.join(os.getcwd(), filename)

    if not os.path.exists(path):
        print(f"File not found: {filename}")
        return

    if os.path.isdir(path):
        print(f"'{filename}' is a directory.")
        return

    try:
        with open(path, "r", encoding="utf-8") as f:
            lines = f.readlines()
            for line in lines[-num_lines:]:
                print(line.rstrip())
    except Exception as e:
        print(f"Error reading file: {e}")
        
@register_command("kill")
def kill_process(args):
    """Terminate a process by its PID (like Unix 'kill')."""
    import psutil

    if not args:
        print("Usage: kill <pid>")
        return

    try:
        pid = int(args[0])
    except ValueError:
        print("Error: PID must be a number.")
        return

    try:
        proc = psutil.Process(pid)
        name = proc.name()
        proc.terminate()
        proc.wait(timeout=3)
        print(f"Process {pid} ({name}) terminated.")
    except psutil.NoSuchProcess:
        print(f"No such process: {pid}")
    except psutil.AccessDenied:
        print(f"Access denied. Try running PyTerm as administrator.")
    except psutil.TimeoutExpired:
        print(f"Process {pid} did not terminate, forcing kill...")
        try:
            proc.kill()
            print(f"Process {pid} forcibly killed.")
        except Exception as e:
            print(f"Failed to kill process: {e}")
    except Exception as e:
        print(f"Error: {e}")
        
        
@register_command("ps")
def ps(args):
    """List running processes.
    - ps            → show all processes
    - ps -p         → show only Python processes
    - ps -s <name>  → search for process name or command (case-insensitive)
    """
    import psutil

    show_python_only = "-p" in args
    search_mode = "-s" in args
    search_term = None

    if search_mode:
        # get the next argument as search term
        try:
            search_index = args.index("-s")
            search_term = args[search_index + 1].lower()
        except IndexError:
            print("Usage: ps -s <name>")
            return

    print(f"{'PID':<8} {'Name':<25} {'CPU%':<6} {'Memory%':<8} {'Command'}")
    print("-" * 100)

    for proc in psutil.process_iter(['pid', 'name', 'cmdline', 'cpu_percent', 'memory_percent']):
        try:
            name = proc.info['name'] or ''
            cmd = " ".join(proc.info['cmdline']) if proc.info['cmdline'] else ''
            cpu = proc.info['cpu_percent']
            mem = proc.info['memory_percent']

            # --- Filtering ---
            if show_python_only and not ("python" in name.lower() or "python" in cmd.lower()):
                continue

            if search_mode:
                if search_term not in name.lower() and search_term not in cmd.lower():
                    continue

            # --- Display ---
            print(f"{proc.info['pid']:<8} {name[:25]:<25} {cpu:<6.1f} {mem:<8.2f} {cmd}")

        except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.ZombieProcess):
            continue
   
@register_command("echo")
def echo(args):
    """Print text to the terminal (like Unix echo)."""
    if not args:
        print()
        return

    # Join everything after echo into a single line
    message = " ".join(args)
    print(message)
   
@register_command("find")
def find(args):
    """Search the entire system for a file or folder name (fast partial match)."""
    import threading

    if not args:
        print("Usage: find <name>")
        return

    search_term = args[0].lower()
    found = []

    print(f"Searching system for '{search_term}'...\n")

    # Choose starting points (root drives or /)
    roots = []
    if os.name == "nt":
        import string
        from ctypes import windll
        bitmask = windll.kernel32.GetLogicalDrives()
        for letter in string.ascii_uppercase:
            if bitmask & 1:
                roots.append(f"{letter}:\\")
            bitmask >>= 1
    else:
        roots = ["/"]

    # Threaded search for speed
    def search_path(root):
        for dirpath, dirnames, filenames in os.walk(root, topdown=True):
            # Ignore very large system folders for speed
            dirnames[:] = [d for d in dirnames if d.lower() not in ("windows", "program files", "programdata", "appdata", "system volume information", "$recycle.bin")]
            for name in filenames + dirnames:
                if search_term in name.lower():
                    full_path = os.path.join(dirpath, name)
                    found.append(full_path)
                    print(full_path)

    threads = []
    for r in roots:
        t = threading.Thread(target=search_path, args=(r,))
        t.daemon = True
        threads.append(t)
        t.start()

    for t in threads:
        t.join()

    if not found:
        print("\nNo matching files or folders found.")
    else:
        print(f"\nFound {len(found)} match(es).")
        
@register_command("watch")
def watch(args):
    """Run a command repeatedly (like Unix 'watch').
    Usage:
      watch -n <seconds> <command>
      watch -n <seconds> -t <count> <command>
    """
    if not args or "-n" not in args:
        print("Usage: watch -n <seconds> [-t <count>] <command>")
        return

    # Default values
    interval = 2
    count = None

    try:
        n_index = args.index("-n")
        interval = float(args[n_index + 1])
    except (ValueError, IndexError):
        print("Error: missing or invalid value for -n <seconds>")
        return

    # Optional -t flag for total runs
    if "-t" in args:
        try:
            t_index = args.index("-t")
            count = int(args[t_index + 1])
        except (ValueError, IndexError):
            print("Error: missing or invalid value for -t <count>")
            return

    # Get command after flags
    try:
        # Strip flags and their arguments
        stripped = []
        skip_next = False
        for i, a in enumerate(args):
            if skip_next:
                skip_next = False
                continue
            if a in ("-n", "-t"):
                skip_next = True
                continue
            stripped.append(a)

        if not stripped:
            print("Error: no command provided to watch.")
            return

        command_line = " ".join(stripped)
    except Exception:
        print("Error parsing command.")
        return

    print(f"Watching command: '{command_line}' every {interval}s" + (f" ({count} times)" if count else " (until stopped)"))
    print("Press Ctrl+C to stop.\n")

    runs = 0
    try:
        while True:
            runs += 1
            os.system("cls" if os.name == "nt" else "clear")
            print(f"--- Run {runs} ---\n")
            handle_command(command_line)

            if count and runs >= count:
                break
            time.sleep(interval)
    except KeyboardInterrupt:
        print("\nStopped watching.")      

@register_command("grep")
def grep(args):
    """Search for text in files (like Unix 'grep').
    Usage:
      grep <pattern> <filename>
      grep -i <pattern> <filename>     # case-insensitive
    """
    if not args:
        print("Usage: grep [-i] <pattern> <filename>")
        return

    ignore_case = False
    pattern = None
    filename = None

    # Parse arguments
    if args[0] == "-i":
        ignore_case = True
        if len(args) < 3:
            print("Usage: grep -i <pattern> <filename>")
            return
        pattern = args[1]
        filename = args[2]
    else:
        if len(args) < 2:
            print("Usage: grep <pattern> <filename>")
            return
        pattern = args[0]
        filename = args[1]

    file_path = os.path.join(os.getcwd(), filename)

    if not os.path.exists(file_path):
        print(f"File not found: {filename}")
        return

    if os.path.isdir(file_path):
        print(f"'{filename}' is a directory.")
        return

    # Perform the search
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()

        matches = 0
        for i, line in enumerate(lines, start=1):
            text = line.rstrip("\n")
            if ignore_case:
                if pattern.lower() in text.lower():
                    print(f"{i:>4}: {text}")
                    matches += 1
            else:
                if pattern in text:
                    print(f"{i:>4}: {text}")
                    matches += 1

        if matches == 0:
            print(f"No matches for '{pattern}' in {filename}.")
    except Exception as e:
        print(f"Error reading file: {e}")

@register_command("sudo")
def sudo_command(args):
    """Run a Windows command as administrator (like Linux sudo)."""
    if os.name != "nt":
        print("sudo is only available on Windows.")
        return

    if not args:
        print("Usage: sudo <windows_command>")
        return

    import ctypes
    import subprocess

    cmd_line = " ".join(args)

    try:
        # Run PowerShell as Administrator
        print(f"Running as administrator: {cmd_line}")
        ctypes.windll.shell32.ShellExecuteW(
            None,
            "runas",
            "powershell.exe",
            f'/c {cmd_line}',
            None,
            1
        )
    except Exception as e:
        print(f"Error running command as admin: {e}")


@register_command("neofetch")
def neofetch(args):
    """Display system info with PyNixShell logo (like Linux neofetch)."""
    import platform
    import psutil
    from datetime import datetime
    import re

    # --- PyNixShell ASCII logo ---
    ascii_logo = r"""
██████╗ ██╗   ██╗
██╔══██╗╚██╗ ██╔╝
██████╔╝ ╚████╔╝ 
██╔═══╝   ╚██╔╝  
██║        ██║   
╚═╝        ╚═╝   
███╗   ██╗██╗██╗  ██╗
████╗  ██║██║╚██╗██╔╝
██╔██╗ ██║██║ ╚███╔╝ 
██║╚██╗██║██║ ██╔██╗ 
██║ ╚████║██║██╔╝ ██╗
╚═╝  ╚═══╝╚═╝╚═╝  ╚═╝
                              The Power of Python + Unix
    """

    # --- System info ---
    uname = platform.uname()
    mem = psutil.virtual_memory()
    disk = psutil.disk_usage('/')
    boot_time = datetime.fromtimestamp(psutil.boot_time()).strftime("%Y-%m-%d %H:%M:%S")

    sysinfo = [
        f"OS: {uname.system} {uname.release}",
        f"Host: {uname.node}",
        f"Kernel: {uname.version}",
        f"CPU: {uname.processor or 'Unknown CPU'}",
        f"Cores: {psutil.cpu_count(logical=False)} ({psutil.cpu_count(logical=True)} threads)",
        f"Memory: {mem.used // (1024**2)}MB / {mem.total // (1024**2)}MB",
        f"Disk: {disk.used // (1024**3)}GB / {disk.total // (1024**3)}GB",
        f"Boot: {boot_time}",
        f"Shell: PyNixShell v1.0",
    ]

    # --- Combine logo and info side-by-side ---
    ascii_lines = ascii_logo.splitlines()
    max_logo_width = max(len(re.sub(r"\x1b\[[0-9;]*m", "", line)) for line in ascii_lines if line.strip())

    print()
    for i in range(max(len(ascii_lines), len(sysinfo))):
        logo_line = ascii_lines[i] if i < len(ascii_lines) else ""
        info_line = sysinfo[i] if i < len(sysinfo) else ""
        print(f"{logo_line.ljust(max_logo_width + 4)}{info_line}")
    print()
    
@register_command("shutdown")
def shutdown_cmd(args):
    """Shut down the system (asks for confirmation)."""
    if not confirm_action("Confirm Shutdown", "Are you sure you want to shut down your computer?"):
        print("Shutdown cancelled.")
        return

    system = platform.system()
    if system == "Windows":
        os.system("shutdown /s /t 0")
    elif system == "Linux":
        os.system("shutdown now")
    elif system == "Darwin":  # macOS
        os.system("osascript -e 'tell app \"System Events\" to shut down'")
    else:
        print("Unsupported OS for shutdown.")
    print("Shutting down...")


@register_command("restart")
def restart_cmd(args):
    """Restart the system (asks for confirmation)."""
    if not confirm_action("Confirm Restart", "Are you sure you want to restart your computer?"):
        print("Restart cancelled.")
        return

    system = platform.system()
    if system == "Windows":
        os.system("shutdown /r /t 0")
    elif system == "Linux":
        os.system("reboot")
    elif system == "Darwin":  # macOS
        os.system("osascript -e 'tell app \"System Events\" to restart'")
    else:
        print("Unsupported OS for restart.")
    print("Restarting...")


@register_command("sleep")
def sleep_cmd(args):
    """Put the computer to sleep immediately."""
    system = platform.system()
    if system == "Windows":
        ctypes.windll.PowrProf.SetSuspendState(0, 0, 0)
    elif system == "Linux":
        os.system("systemctl suspend")
    elif system == "Darwin":  # macOS
        os.system("pmset sleepnow")
    else:
        print("Unsupported OS for sleep.")
    print("💤 Going to sleep...")    

@register_command("help")
def commands_cmd(args):
    """Displays all commands, search results, or detailed info from commands.json."""
    import json

    file_path = "commands.json"

    if not os.path.exists(file_path):
        print("❌ commands.json not found.")
        return

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        print(f"⚠️ Error reading {file_path}: {e}")
        return

    # Convert list into lookup dict
    all_cmds = {cmd["name"].lower(): cmd for cmd in data}

    # --- commands search <keyword> ---
    if len(args) >= 2 and args[0].lower() == "search":
        keyword = " ".join(args[1:]).lower()
        print(f"🔍 Searching for '{keyword}'...\n")
        found = False
        for cmd in data:
            if (
                keyword in cmd["name"].lower()
                or keyword in cmd.get("desc", "").lower()
                or keyword in cmd.get("definition", "").lower()
                or keyword in cmd.get("category", "").lower()
            ):
                print(f"  {cmd['name']:<15} {cmd['desc']}")
                found = True
        if not found:
            print("No commands matched your search.")
        return

    # --- commands <command> ---
    if len(args) == 1:
        cmd_name = args[0].lower()
        if cmd_name in all_cmds:
            cmd_info = all_cmds[cmd_name]
            print(f"📘 {cmd_info['name']} ({cmd_info.get('category', 'other')})\n")
            print(f"Description: {cmd_info['desc']}")
            if "definition" in cmd_info:
                print(f"\nDetails:\n{cmd_info['definition']}")
            return
        else:
            print(f"Command '{cmd_name}' not found.")
            return

    # --- Default: list all commands grouped by category ---
    print("📘 Available Commands in PyNixShell:\n")
    categories = {}
    for cmd in data:
        cat = cmd.get("category", "other").capitalize()
        categories.setdefault(cat, []).append(cmd)

    for cat, cmds in sorted(categories.items()):
        print(f"# --- {cat.upper()} ---")
        for cmd in sorted(cmds, key=lambda c: c["name"]):
            print(f"  {cmd['name']:<15} {cmd['desc']}")
        print()

    print("\n💡 Use 'help <name>' for details or commands to view all commands with no info.\nPrefix any Windows command with 'win' (e.g. win dir, win cls, win ping, win cd for basic windows commands)")

@register_command("view")
def view_folder(args):
    """View contents of a folder or archive without entering it.
    Usage:
      view                → shows current directory contents
      view <path>         → shows contents of specified folder
      view -z <file.zip>  → lists contents of a ZIP archive
      view -t <file.tar>  → lists contents of a TAR archive
    """

    def human_readable(size_bytes):
        """Convert bytes to human-readable format."""
        if size_bytes == 0:
            return "0 B"
        size_name = ("B", "KB", "MB", "GB", "TB")
        i = int(math.floor(math.log(size_bytes, 1024)))
        p = math.pow(1024, i)
        s = round(size_bytes / p, 2)
        return f"{s} {size_name[i]}"

    # --- Handle ZIP mode ---
    if args and args[0] == "-z":
        if len(args) < 2:
            print("Usage: view -z <file.zip>")
            return
        zip_path = args[1]
        if not os.path.exists(zip_path):
            print(f"❌ ZIP file not found: {zip_path}")
            return
        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                info_list = zip_ref.infolist()
                print(f"\n🗜️ Contents of {zip_path}:\n")
                for info in info_list:
                    size = human_readable(info.file_size)
                    print(f"  📄 {info.filename:<50} {size:>10}")
        except zipfile.BadZipFile:
            print(f"❌ Error: '{zip_path}' is not a valid ZIP archive.")
        return

    # --- Handle TAR mode ---
    if args and args[0] == "-t":
        if len(args) < 2:
            print("Usage: view -t <file.tar>")
            return
        tar_path = args[1]
        if not os.path.exists(tar_path):
            print(f"❌ TAR file not found: {tar_path}")
            return
        try:
            with tarfile.open(tar_path, "r") as tar_ref:
                print(f"\n📦 Contents of {tar_path}:\n")
                for member in tar_ref.getmembers():
                    size = human_readable(member.size)
                    prefix = "📁" if member.isdir() else "📄"
                    print(f"  {prefix} {member.name:<50} {size:>10}")
        except tarfile.TarError:
            print(f"❌ Error: '{tar_path}' is not a valid TAR archive.")
        return

    # --- Normal folder viewing ---
    path = os.getcwd() if not args else args[0]

    if not os.path.exists(path):
        print(f"❌ Path not found: {path}")
        return

    if not os.path.isdir(path):
        print(f"'{path}' is not a directory.")
        return

    items = sorted(os.listdir(path))
    if not items:
        print("(empty folder)")
        return

    print(f"\n📂 Contents of {path}:\n")
    for item in items:
        item_path = os.path.join(path, item)
        if os.path.isdir(item_path):
            print(f"  📁 {item}/")
        else:
            size = os.path.getsize(item_path)
            print(f"  📄 {item:<30} {human_readable(size):>10}")
            
@register_command("registercommand")
def registercommand(args):
    """Explains how to register and add an external command file."""
    print("""
📘 How to Create and Register a Custom Command in PyNixShell
─────────────────────────────────────────────────────────────

You can create your own commands as Python files and add them permanently to PyNixShell.

1️⃣  Create a new `.py` file (example: `hello.py`) with this template:

    import os

    def register_command(name):
        from __main__ import registered_commands
        def wrapper(func):
            registered_commands[name] = func
            return func
        return wrapper


    @register_command("hello")
    def hello(args):
        \"\"\"Prints a simple greeting.\"\"\"
        print("Hello from your custom command!")

─────────────────────────────────────────────────────────────

2️⃣  Add your command to PyNixShell:
    
    add hello.py

   ✅ This copies it into `/commands/added/` and loads it instantly.
   ✅ It will also auto-load next time you start PyNixShell.

─────────────────────────────────────────────────────────────

3️⃣  Test it:
    
    hello

─────────────────────────────────────────────────────────────

💡 Tip: You can organize your permanent commands inside:
    /commands/added/
Each .py file here is loaded automatically at startup.
""")
    
@register_command("add")
def add_command(args):
    """Add (copy) a Python command file into /commands for permanent use.
    
    Usage:
        add file.py
    Copies file.py into /commands/ and loads it immediately.
    """
    if not args:
        print("Usage: add <file.py>")
        return

    src = args[0]

    # Validate
    if not os.path.exists(src):
        print(f"❌ File not found: {src}")
        return

    if not src.endswith(".py"):
        print("❌ Only .py files are supported.")
        return

    # Ensure /commands folder exists
    commands_dir = os.path.join(os.getcwd(), "commands")
    os.makedirs(commands_dir, exist_ok=True)

    dest = os.path.join(commands_dir, os.path.basename(src))

    # Check for duplicates
    if os.path.exists(dest):
        choice = input(f"⚠️ Command '{os.path.basename(src)}' already exists. Overwrite? (y/n): ").strip().lower()
        if choice != "y":
            print("❌ Operation cancelled.")
            return

    try:
        # Copy file
        shutil.copy2(src, dest)
        print(f"✅ Copied '{os.path.basename(src)}' to /commands/")

        # Load immediately
        spec = importlib.util.spec_from_file_location(os.path.basename(dest), dest)
        module = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(module)

        print("🔁 Command loaded and ready to use.")
    except Exception as e:
        print(f"❌ Failed to add command: {e}")    
    
@register_command("commands")
def list_registered_commands(args):
    """Lists all currently registered commands in PyNixShell."""
    print("📘 Registered Commands in PyNixShell:\n")

    if not registered_commands:
        print("⚠️ No commands registered yet.")
        return

    # Sort alphabetically for readability
    for cmd in sorted(registered_commands.keys()):
        print(f"  {cmd}")

    print(f"\n🧩 Total: {len(registered_commands)} commands loaded.")
    print("💡 Tip: Use 'help <command>' for more details.")
    
@register_command("rsync")
def rsync(args):
    """Synchronize all matching files or a specific target file with the source.
    Usage:
      rsync <path/to/source.ext>             → sync all files on system matching that name
      rsync <path/to/source.ext> -t <target> → sync only to a specific file or folder
    """
    import hashlib
    import threading

    # --- Argument validation ---
    if not args:
        print("Usage: rsync <source> [-t <target>]")
        return

    source_path = os.path.abspath(args[0])
    if not os.path.exists(source_path):
        print(f"❌ Source file not found: {source_path}")
        return

    # Optional target flag
    target_path = None
    if "-t" in args:
        try:
            target_path = os.path.abspath(args[args.index("-t") + 1])
        except IndexError:
            print("⚠️ Missing argument for '-t'")
            return

    # --- Hash helper ---
    def file_hash(path):
        h = hashlib.md5()
        try:
            with open(path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    h.update(chunk)
            return h.hexdigest()
        except Exception:
            return None

    source_hash = file_hash(source_path)
    filename = os.path.basename(source_path)
    size = os.path.getsize(source_path)

    # --- If a specific target was provided ---
    if target_path:
        # Handle: folder, existing file, or new file
        if os.path.isdir(target_path):
            dest_file = os.path.join(target_path, filename)
        else:
            dest_file = target_path
            os.makedirs(os.path.dirname(dest_file), exist_ok=True)

        try:
            shutil.copy2(source_path, dest_file)
            print(f"✅ Synchronized {source_path} → {dest_file}")
        except Exception as e:
            print(f"❌ Failed to sync: {e}")
        return

    # --- Otherwise, perform system-wide sync ---
    print(f"🔍 Scanning for other '{filename}' files across system...\n")

    matches = []

    # Collect root directories
    roots = []
    if os.name == "nt":
        import string
        from ctypes import windll
        bitmask = windll.kernel32.GetLogicalDrives()
        for letter in string.ascii_uppercase:
            if bitmask & 1:
                roots.append(f"{letter}:\\")
            bitmask >>= 1
    else:
        roots = ["/"]

    # Search
    def search_path(root):
        for dirpath, dirnames, filenames in os.walk(root, topdown=True):
            # Skip system folders for speed/safety
            dirnames[:] = [d for d in dirnames if d.lower() not in (
                "windows", "program files", "programdata", "appdata",
                "system volume information", "$recycle.bin")]
            for f in filenames:
                if f.lower() == filename.lower():
                    full = os.path.join(dirpath, f)
                    if os.path.abspath(full) != source_path:
                        matches.append(full)
                        print(full)

    # Threaded search
    threads = []
    for root in roots:
        t = threading.Thread(target=search_path, args=(root,))
        t.daemon = True
        threads.append(t)
        t.start()
    for t in threads:
        t.join()

    print(f"\n📁 Found {len(matches)} matching file(s).")

    if not matches:
        return

    confirm = input(f"\n⚠️ Overwrite all {len(matches)} file(s) with {source_path}? (y/n): ").strip().lower()
    if confirm != "y":
        print("❌ Sync cancelled.")
        return

    replaced = 0
    for target in matches:
        try:
            if os.path.getsize(target) == size and file_hash(target) == source_hash:
                continue
            shutil.copy2(source_path, target)
            replaced += 1
            print(f"✅ Synced: {target}")
        except Exception as e:
            print(f"⚠️ Failed to sync {target}: {e}")

    print(f"\n✅ Sync complete. {replaced}/{len(matches)} file(s) updated.")
    
@register_command("option")
def contextmenu(args):
    """Interactive right-click style menu for a file or folder."""
    if not args:
        print("Usage: contextmenu <filename or folder>")
        return

    target = args[0]
    path = os.path.join(os.getcwd(), target)

    if not os.path.exists(path):
        print(f"❌ Not found: {target}")
        return

    is_dir = os.path.isdir(path)
    ext = os.path.splitext(path)[1].lower()

    # --- Define menu items dynamically ---
    menu_items = ["Open", "Copy", "Cut", "Delete", "Rename",]

    if is_dir:
        menu_items += [
            "Open in Explorer" if os.name == "nt" else "Open in Finder",
            "New File",
            "New Folder",
            "Compress to ZIP",
        ]
    else:
        if ext in (".py", ".txt", ".json", ".html", ".js"):
            menu_items += ["Edit", "View"]
        if ext in (".exe", ".bat", ".cmd", ".sh", ".app"):
            menu_items += ["Run as Administrator" if os.name == "nt" else "Run"]
        if ext in (".zip", ".tar", ".gz"):
            menu_items += ["Extract Here"]


    if os.name == "nt":
        menu_items += ["Send to → Desktop (Shortcut)", "Send to → Documents"]
    else:
        menu_items += ["Copy Path", "Open Terminal Here"]

    # --- Curses UI ---
    def menu(stdscr):
        curses.curs_set(0)
        stdscr.clear()
        h, w = stdscr.getmaxyx()

        current_idx = 0

        while True:
            stdscr.clear()
            title = f"📂 Context Menu: {os.path.basename(path)}"
            stdscr.addstr(0, 0, title[:w-1], curses.A_BOLD | curses.A_UNDERLINE)

            for i, item in enumerate(menu_items):
                x = 2
                y = i + 2
                if i == current_idx:
                    stdscr.addstr(y, x, f"> {item}", curses.A_REVERSE)
                else:
                    stdscr.addstr(y, x, f"  {item}")

            stdscr.refresh()
            key = stdscr.getch()

            if key in (curses.KEY_UP, ord('k')):
                current_idx = (current_idx - 1) % len(menu_items)
            elif key in (curses.KEY_DOWN, ord('j')):
                current_idx = (current_idx + 1) % len(menu_items)
            elif key in (10, 13):  # Enter
                chosen = menu_items[current_idx]
                stdscr.clear()
                stdscr.addstr(0, 0, f"✅ Selected: {chosen}")
                stdscr.refresh()
                curses.napms(500)
                break
            elif key in (27, ord('q')):  # ESC or Q
                chosen = None
                break

        return chosen

    try:
        selected = curses.wrapper(menu)
        if not selected:
            print("❌ Cancelled.")
            return

        # --- Perform actions ---
        print(f"\nSelected: {selected}")
        if selected == "Open":
            if is_dir:
                os.startfile(path) if os.name == "nt" else subprocess.call(["open", path])
            else:
                os.startfile(path) if os.name == "nt" else subprocess.call(["open", path])
        elif selected == "Edit":
            handle_command(f"nano {target}")
            
        elif selected == "Delete":
            handle_command(f"rm {target}")
        elif selected == "Copy":
            handle_command(f"cp {target}")
        elif selected == "Paste":
            handle_command("paste")
        elif selected == "Rename":
            new_name = input("Enter new name: ")
            if new_name:
                handle_command(f"move {target} {new_name}")
        elif selected == "New Folder":
            handle_command("mkdir NewFolder")
        elif selected == "New File":
            handle_command("touch newfile.txt")
        elif selected == "Open Terminal Here":
            os.chdir(path)
            print(f"Opened terminal in: {path}")
        else:
            print(f"(No direct command mapped for '{selected}')")

    except Exception as e:
        print(f"Error: {e}")    

@register_command("history")
def show_history(args):
    """Display a list of previously run commands.
    Options:
      history           → show all commands
      history -n <num>  → show last <num> commands
      history -c        → clear history
    """
    global command_history

    if not command_history:
        print("No commands in history.")
        return

    # Clear history
    if args and args[0].lower() in ("-c", "--clear"):
        command_history.clear()
        print("🧹 Command history cleared.")
        return

    # Show last N commands
    if args and args[0] == "-n":
        if len(args) < 2 or not args[1].isdigit():
            print("Usage: history -n <number>")
            return
        n = int(args[1])
        entries = command_history[-n:]
        print(f"\n📜 Last {n} Commands:")
    else:
        entries = command_history
        print("\n📜 Command History:")

    for i, cmd in enumerate(entries, start=len(command_history) - len(entries) + 1):
        print(f"{i:>3}. {cmd}")

@register_command("ip")
def ip_info(args):
    """Displays local or public IP information (like ipconfig/ifconfig)."""
    try:
        # --- Simple Mode: show system IPv4(s) ---
        if args and args[0] == "-s":
            print("📡 Active IPv4 Addresses:\n")
            found = False
            for interface, addrs in psutil.net_if_addrs().items():
                for addr in addrs:
                    if addr.family == socket.AF_INET:
                        print(f"  {interface:<20} → {addr.address}")
                        found = True
            if not found:
                print("⚠️  No active IPv4 addresses found.")
            return

        # --- Default mode: show hostname, local, and public IP ---
        hostname = socket.gethostname()
        local_ip = socket.gethostbyname(hostname)
        print(f"📡 Hostname: {hostname}")
        print(f"🌐 Local IP: {local_ip}")

        # --- Try to get public IP ---
        import urllib.request
        urls = [
            "https://api.ipify.org",
            "https://icanhazip.com",
            "https://ident.me"
        ]
        external_ip = None
        for url in urls:
            try:
                with urllib.request.urlopen(url, timeout=3) as response:
                    external_ip = response.read().decode().strip()
                    if external_ip:
                        break
            except Exception:
                continue

        if external_ip:
            print(f"🌎 Public IP: {external_ip}")
        else:
            print("🌎 Public IP: (unavailable — no network or blocked)")

    except Exception as e:
        print(f"Error retrieving IP info: {e}")
        
@register_command("donut")
def donut(args):
    handle_command("clear")
    print(" Press \"ESC\" key to stop")
    """Smooth spinning ASCII donut animation (press ESC to stop)."""
    global DonutRunning
    width, height = 60, 24
    chars = list(".,-~:;=!*#$@")
    A, B = 1, 1
    DonutRunning = True

    print("🍩 Spinning donut! (Press ESC to stop)\n")

    def get_key_pressed():
        """Cross-platform ESC key detection."""
        if os.name == "nt":  # Windows
            if msvcrt.kbhit():
                key = msvcrt.getch()
                return key == b'\x1b'  # ESC
        else:  # macOS/Linux
            dr, dw, de = select.select([sys.stdin], [], [], 0)
            if dr:
                key = sys.stdin.read(1)
                return key == '\x1b'
        return False

    def run_donut():
        global DonutRunning
        nonlocal A, B
        try:
            while DonutRunning:
                # Build frame without clearing terminal
                b = [" "] * (width * height)
                z = [0] * (width * height)

                for j in [x * 0.07 for x in range(int(6.28 / 0.07))]:
                    for i in [x * 0.02 for x in range(int(6.28 / 0.02))]:
                        c, l = math.sin(i), math.cos(i)
                        d, f = math.cos(j), math.sin(j)
                        g, e = math.cos(A), math.sin(A)
                        h = d + 2
                        D = 1 / (c * h * e + f * g + 5)
                        m, n = math.cos(B), math.sin(B)
                        t = c * h * g - f * e

                        x = int(width / 2 + 30 * D * (l * h * m - t * n))
                        y = int(height / 2 + 15 * D * (l * h * n + t * m))
                        o = x + width * y
                        N = int(8 * ((f * e - c * d * g) * m - c * d * e - f * g - l * d * n))

                        if 0 <= y < height and 0 <= x < width and D > z[o]:
                            z[o] = D
                            b[o] = chars[max(0, min(len(chars) - 1, N))]

                frame = ""
                for i in range(len(b)):
                    frame += b[i]
                    if i % width == 0:
                        frame += "\n"

                # Move cursor to top (no flicker)
                sys.stdout.write("\033[H")
                sys.stdout.write(frame)
                sys.stdout.flush()

                A += 0.07
                B += 0.03
                time.sleep(0.05)

                if get_key_pressed():
                    DonutRunning = False
                    break

        except Exception as e:
            print(f"\nError: {e}")
        finally:
            sys.stdout.write("\033[0m\n")
            sys.stdout.flush()
            print("\n🍩 Donut stopped.")

    # On Unix, set terminal to raw mode so ESC key is captured immediately
    if os.name != "nt":
        fd = sys.stdin.fileno()
        old_settings = termios.tcgetattr(fd)
        tty.setcbreak(fd)

    threading.Thread(target=run_donut, daemon=True).start()

    # Wait for donut to stop (so user can ESC anytime)
    while DonutRunning:
        time.sleep(0.1)

    # Restore terminal state on Unix
    if os.name != "nt":
        termios.tcsetattr(fd, termios.TCSADRAIN, old_settings)
        
@register_command("zip")
def zip_command(args):
    """Create a ZIP archive using the system’s native zip tools."""
    if len(args) < 2:
        print("Usage: zip <source_folder> <destination_name.zip>")
        return

    source = args[0]
    dest = args[1]

    # Ensure .zip extension
    if not dest.lower().endswith(".zip"):
        dest += ".zip"

    # Verify source
    if not os.path.exists(source):
        print(f"Source not found: {source}")
        return

    # --- Windows: use PowerShell Compress-Archive ---
    if os.name == "nt":
        try:
            cmd = [
                "powershell",
                "-Command",
                f'Compress-Archive -Path "{source}" -DestinationPath "{dest}" -Force'
            ]
            subprocess.run(cmd, check=True)
            print(f"✅ Created ZIP archive: {dest}")
            return
        except subprocess.CalledProcessError as e:
            print(f"Error: {e}")
            return

    # --- macOS/Linux: use system 'zip' command if available ---
    elif shutil.which("zip"):
        try:
            subprocess.run(["zip", "-r", dest, source], check=True)
            print(f"✅ Created ZIP archive: {dest}")
        except subprocess.CalledProcessError as e:
            print(f"Error: {e}")
    else:
        # --- Fallback: use Python’s built-in make_archive ---
        try:
            shutil.make_archive(dest.replace(".zip", ""), "zip", source)
            print(f"✅ Created ZIP archive (fallback): {dest}")
        except Exception as e:
            print(f"Error creating ZIP archive: {e}")     


@register_command("tar")
def tar_command(args):
    """Create a TAR archive using the native system tar command."""
    if len(args) < 2:
        print("Usage: tar <source_folder> <destination.tar>")
        return

    source = args[0]
    dest = args[1]

    # Ensure .tar extension
    if not dest.lower().endswith(".tar"):
        dest += ".tar"

    # Verify the source exists
    if not os.path.exists(source):
        print(f"❌ Source not found: {source}")
        return

    # --- Use native tar if available ---
    if shutil.which("tar"):
        try:
            subprocess.run(["tar", "-cvf", dest, source], check=True)
            print(f"✅ Created TAR archive: {dest}")
        except subprocess.CalledProcessError as e:
            print(f"Error creating TAR archive: {e}")
        return

    # --- Fallback using Python tarfile module ---
    import tarfile
    try:
        with tarfile.open(dest, "w") as tar:
            tar.add(source, arcname=os.path.basename(source))
        print(f"✅ Created TAR archive (fallback): {dest}")
    except Exception as e:
        print(f"Error creating TAR archive: {e}")        

@register_command("alias")
def alias_cmd(args):
    """Create or list command aliases.
    Usage:
      alias                       → list all aliases
      alias name='command'        → create alias
    Example:
      alias ll='ls -l'
    """
    global _aliases

    # --- List all aliases ---
    if not args:
        if not _aliases:
            print("(no aliases defined)")
            return
        print("\n📘 Aliases:\n")
        for name, cmd in _aliases.items():
            print(f"  {name} → {cmd}")
        return

    # --- Parse and create alias ---
    line = " ".join(args).strip()
    if "=" not in line:
        print("Usage: alias name='command'")
        return

    name, cmd = line.split("=", 1)
    name = name.strip()
    cmd = cmd.strip().strip("'\"")

    # --- Register alias ---
    _aliases[name] = cmd
    registered_commands[name] = lambda a, c=cmd: handle_command(c)

    # --- Save it ---
    save_aliases()
    print(f"✅ Alias saved: {name} → {cmd}")

@register_command("unalias")
def unalias_cmd(args):
    """Remove an alias.
    Usage:
      unalias <name>
      unalias -a       → remove all aliases
    """
    global _aliases

    if not args:
        print("Usage: unalias <name> or unalias -a")
        return

    if args[0] == "-a":
        count = len(_aliases)
        _aliases.clear()
        save_aliases()
        print(f"🗑️  Removed all {count} aliases.")
        return

    name = args[0]
    if name not in _aliases:
        print(f"❌ Alias not found: {name}")
        return

    # Remove from memory and disk
    del _aliases[name]
    if name in registered_commands:
        del registered_commands[name]
    save_aliases()
    print(f"🗑️  Alias removed: {name}")
    
@register_command("date")
def date_cmd(args):
    """Display the current date and time."""
    now = datetime.datetime.now()

    # Default format (similar to Unix `date`)
    formatted = now.strftime("%a %b %d %H:%M:%S %Y")

    # Support a flag for ISO format if desired
    if args and args[0] == "-i":
        formatted = now.isoformat()

    print(f"📅 {formatted}")    
    
@register_command("csync")
def csync_cmd(args):
    """
    Synchronize and reload all command files in /commands (recursively).
    Deletes stale modules, reloads updated ones, and imports new ones.
    """
    base_dir = os.path.join(os.getcwd(), "commands")
    os.makedirs(base_dir, exist_ok=True)

    print("🔄 Starting command sync...\n")

    # 1️⃣ Collect all .py files inside /commands recursively
    all_command_files = []
    for root, _, files in os.walk(base_dir):
        for file in files:
            if file.endswith(".py"):
                file_path = os.path.join(root, file)
                all_command_files.append(file_path)

    if not all_command_files:
        print("⚠️ No command files found in /commands.")
        return

    # 2️⃣ Identify currently loaded external modules
    loaded_modules = {
        name: module for name, module in sys.modules.items()
        if name.startswith("cmd_")
    }

    # 3️⃣ Remove orphaned (deleted) modules
    to_remove = []
    for mod_name, module in loaded_modules.items():
        mod_file = getattr(module, "__file__", None)
        if not mod_file or not os.path.exists(mod_file):
            to_remove.append(mod_name)

    for mod_name in to_remove:
        del sys.modules[mod_name]
        print(f"🗑️  Removed stale module: {mod_name}")

    # 4️⃣ Reload all valid command files
    import __main__
    loaded_count = 0
    failed = []

    for file_path in all_command_files:
        try:
            rel_path = os.path.relpath(file_path, base_dir)
            hash_suffix = hashlib.md5(rel_path.encode()).hexdigest()[:6]
            mod_name = f"cmd_{os.path.splitext(os.path.basename(file_path))[0]}_{hash_suffix}"

            # Remove old version if it exists
            if mod_name in sys.modules:
                del sys.modules[mod_name]

            # Load new module
            spec = importlib.util.spec_from_file_location(mod_name, file_path)
            if not spec or not spec.loader:
                print(f"⚠️ Skipping invalid module: {file_path}")
                continue

            module = importlib.util.module_from_spec(spec)
            module.__dict__["registered_commands"] = __main__.registered_commands
            module.__dict__["handle_command"] = __main__.handle_command
            spec.loader.exec_module(module)
            sys.modules[mod_name] = module

            loaded_count += 1
            print(f"✅ Loaded: {os.path.basename(file_path)}")

        except Exception as e:
            failed.append((file_path, str(e)))
            traceback.print_exc()

    # 5️⃣ Summary

    print(f"\n✅ use command rfpt to refresh {loaded_count} command file(s).")
    if failed:
        print("⚠️ Failed to load:")
        for name, err in failed:
            print(f"  - {name}: {err}")

    print("\n🔁 All commands are now synced and up to date.")
    
@register_command("lsc")
def lsc_cmd(args):
    """
    Lists external command files and their registered commands.
    Usage:
      lsc         → shows all active loaded commands with file paths
      lsc -c      → scans /commands folder for @register_command("...") definitions
    """
    base_dir = os.path.join(os.getcwd(), "commands")

    if not os.path.exists(base_dir):
        print("⚠️ No /commands directory found.")
        return

    # --- Normal mode: list active loaded commands in memory ---
    if not args or args[0] != "-c":
        print("📦 Active External Commands (from memory):\n")
        found_any = False

        for cmd_name, func in registered_commands.items():
            try:
                file_path = inspect.getsourcefile(func)
                if file_path and os.path.commonpath([file_path, base_dir]) == base_dir:
                    found_any = True
                    abs_path = os.path.abspath(file_path)
                    print(f"🧩 {cmd_name:<20} → {abs_path}")
            except Exception:
                continue

        if not found_any:
            print("⚠️ No external commands currently loaded.")
        return

    # --- -c mode: scan files directly for register_command() calls ---
    print("🔍 Scanning /commands for external command definitions...\n")
    found_any = False
    pattern = re.compile(r'@register_command\((["\'])(.*?)\1\)')

    for root, _, files in os.walk(base_dir):
        for file in files:
            if not file.endswith(".py"):
                continue
            file_path = os.path.join(root, file)

            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()

                matches = [m[1] for m in pattern.findall(content)]
                if matches:
                    found_any = True
                    abs_path = os.path.abspath(file_path)
                    print(f"📄 {abs_path}")
                    for cmd in matches:
                        print(f"   └─ 🧩 {cmd}")
            except Exception as e:
                print(f"⚠️ Error reading {file}: {e}")

    if not found_any:
        print("⚠️ No @register_command definitions found in /commands.")
        
@register_command("mod")
def mod_cmd(args):
    """
    Run a command after a time delay.
    Usage:
      mod -t <seconds> <command>
    Example:
      mod -t 5 echo Hello world
    """
    if not args or len(args) < 3 or args[0] != "-t":
        print("Usage: mod -t <seconds> <command>")
        return

    try:
        delay = float(args[1])
    except ValueError:
        print("Error: time must be a number (seconds).")
        return

    command_line = " ".join(args[2:])
    print(f"⏳ Waiting {delay} seconds before running: {command_line}")

    def delayed_execute():
        time.sleep(delay)
        print(f"\n▶ Running delayed command: {command_line}")
        handle_command(command_line)

    # Run asynchronously so terminal isn’t blocked
    threading.Thread(target=delayed_execute, daemon=True).start()        
    
@register_command("newline")
def newline_cmd(args):
    """
    Inserts a blank input line in the terminal (like pressing Enter).
    Usage:
      newline
    """
    print("")  # Just print an empty line
    
@register_command("queue")
def queue_cmd(args):
    """Add a command to autoexec.json to run automatically next startup."""
    if not args:
        print("Usage: queue <command>")
        return

    cmd_line = " ".join(args)
    queue = []
    if os.path.exists(AUTOEXEC_FILE):
        try:
            with open(AUTOEXEC_FILE, "r", encoding="utf-8") as f:
                queue = json.load(f)
        except Exception:
            queue = []

    queue.append(cmd_line)
    with open(AUTOEXEC_FILE, "w", encoding="utf-8") as f:
        json.dump(queue, f, indent=2)
    print(f"✅ Queued for next startup: {cmd_line}")    
    
    
@register_command("unzip")
def unzip_cmd(args):
    """Unzip a .zip file into the current directory."""
    import zipfile

    if not args:
        print("Usage: unzip <file.zip>")
        return

    zip_path = args[0]
    if not os.path.exists(zip_path):
        print(f"❌ File not found: {zip_path}")
        return

    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(os.getcwd())
        print(f"✅ Extracted '{zip_path}' into {os.getcwd()}")
    except zipfile.BadZipFile:
        print("⚠️ Invalid ZIP file or corrupted archive.")
    except Exception as e:
        print(f"⚠️ Error extracting ZIP: {e}")

@register_command("dtar")
def dtar_cmd(args):
    """Extract a .tar or .tar.gz archive into the current directory."""
    import tarfile

    if not args:
        print("Usage: dtar <file.tar> or dtar <file.tar.gz>")
        return

    tar_path = args[0]
    if not os.path.exists(tar_path):
        print(f"❌ File not found: {tar_path}")
        return

    try:
        with tarfile.open(tar_path, "r:*") as tar_ref:
            tar_ref.extractall(os.getcwd())
        print(f"✅ Extracted '{tar_path}' into {os.getcwd()}")
    except tarfile.ReadError:
        print("⚠️ Invalid or corrupted TAR archive.")
    except Exception as e:
        print(f"⚠️ Error extracting TAR: {e}")
    
    
@register_command("ungit")
def ungit_cmd(args):
    """Removes Git repository tracking (.git folder) from a directory (Windows-safe)."""
    import shutil
    import stat

    target_dir = args[0] if args else os.getcwd()

    if not os.path.exists(target_dir):
        print(f"❌ Directory not found: {target_dir}")
        return

    git_dir = os.path.join(target_dir, ".git")

    if not os.path.exists(git_dir):
        print("⚠️ No .git directory found — this folder is not a Git repository.")
        return

    def on_rm_error(func, path, exc_info):
        """Handle permission errors by making file writable and retrying."""
        try:
            os.chmod(path, stat.S_IWRITE)  # Remove read-only flag
            func(path)
        except Exception:
            pass  # Ignore if still fails (e.g. locked by system)

    try:
        # Make sure all files are writable before removal
        for root, dirs, files in os.walk(git_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                try:
                    os.chmod(fpath, stat.S_IWRITE)
                except Exception:
                    pass

        # Attempt removal (with error recovery)
        shutil.rmtree(git_dir, onerror=on_rm_error)
        print(f"✅ Successfully removed Git tracking from: {target_dir}")

    except Exception as e:
        print(f"⚠️ Failed to remove Git data: {e}")
        
@register_command("ssh")
def ssh_cmd(args):
    """
    LAN-based peer system for chat and file sharing.
    Usage:
      ssh host   → start server and connect self as client
      ssh join   → discover and join another host
    """
    import socket, threading, struct, os, time, zipfile, io, shutil

    PORT = 5000
    DISCOVERY_PORT = 5001
    BROADCAST_INTERVAL = 2
    clients = {}
    SHARED_DIR = "shared"
    os.makedirs(SHARED_DIR, exist_ok=True)

    def send_packet(sock, data: bytes):
        sock.sendall(struct.pack("!I", len(data)) + data)

    def recv_packet(sock):
        header = sock.recv(4)
        if not header:
            return None
        msg_len = struct.unpack("!I", header)[0]
        data = b""
        while len(data) < msg_len:
            packet = sock.recv(msg_len - len(data))
            if not packet:
                return None
            data += packet
        return data

    def broadcast(sender_conn, message: str):
        for conn in list(clients.keys()):
            if conn != sender_conn:
                try:
                    send_packet(conn, message.encode())
                except:
                    conn.close()
                    clients.pop(conn, None)

    def handle_client(conn, addr):
        try:
            name = recv_packet(conn).decode()
            clients[conn] = name
            print(f"[+] {name} joined from {addr[0]}")
            broadcast(conn, f"SERVER> {name} joined the chat.")

            while True:
                data = recv_packet(conn)
                if not data:
                    break

                # --- handle file / chat commands (same as your script) ---
                if data.startswith(b"UPLOAD:"):
                    _, filename = data.split(b":", 1)
                    filename = filename.decode()
                    file_data = recv_packet(conn)
                    path = os.path.join(SHARED_DIR, filename)
                    with open(path, "wb") as f:
                        f.write(file_data)
                    send_packet(conn, f"SERVER> File '{filename}' uploaded.".encode())

                elif data.startswith(b"PULL:"):
                    _, filename = data.split(b":", 1)
                    filename = filename.decode()
                    path = os.path.join(SHARED_DIR, filename)
                    if os.path.exists(path):
                        with open(path, "rb") as f:
                            file_data = f.read()
                        send_packet(conn, f"FILE:{filename}".encode())
                        send_packet(conn, file_data)
                    else:
                        send_packet(conn, f"SERVER> File '{filename}' not found.".encode())

                elif data.startswith(b"LISTFILES"):
                    files = os.listdir(SHARED_DIR)
                    file_list = "\n".join(files) if files else "(no shared files yet)"
                    send_packet(conn, f"SERVER> Shared files:\n{file_list}".encode())

                elif data.startswith(b"RMV:"):
                    _, target = data.split(b":", 1)
                    target = target.decode().strip()
                    target_path = os.path.join(SHARED_DIR, target)
                    if os.path.exists(target_path):
                        try:
                            if os.path.isdir(target_path):
                                shutil.rmtree(target_path)
                                send_packet(conn, f"SERVER> Directory '{target}' removed.".encode())
                            else:
                                os.remove(target_path)
                                send_packet(conn, f"SERVER> File '{target}' removed.".encode())
                        except Exception as e:
                            send_packet(conn, f"SERVER> Error removing '{target}': {e}".encode())
                    else:
                        send_packet(conn, f"SERVER> '{target}' not found.".encode())

                elif data.startswith(b"HELP"):
                    help_text = (
                        "SERVER> Commands:\n"
                        "  help                - show this help\n"
                        "  listfiles           - list shared files\n"
                        "  push <file>         - upload a file to host\n"
                        "  pushdir <folder>    - upload a folder as zip\n"
                        "  pull <file>         - download a file\n"
                        "  pulldir <folder>    - download folder as zip\n"
                        "  rmv <name>          - remove file/folder\n"
                        "  exit                - leave session"
                    )
                    send_packet(conn, help_text.encode())

                else:
                    msg = data.decode(errors="ignore")
                    print(f"{name}> {msg}")
                    broadcast(conn, f"{name}> {msg}")

        except Exception as e:
            print(f"[x] Error with {addr}: {e}")
        finally:
            cname = clients.pop(conn, "Unknown")
            conn.close()
            broadcast(None, f"SERVER> {cname} disconnected.")

    def get_local_ip():
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        try:
            s.connect(("8.8.8.8", 80))
            ip = s.getsockname()[0]
        except Exception:
            ip = "127.0.0.1"
        finally:
            s.close()
        return ip

    def broadcast_host_info(ip):
        sock = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        sock.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
        name = os.getenv("COMPUTERNAME", socket.gethostname())
        while True:
            msg = f"HOST:{name}:{ip}".encode()
            sock.sendto(msg, ("<broadcast>", DISCOVERY_PORT))
            time.sleep(BROADCAST_INTERVAL)

    def start_server():
        host_ip = get_local_ip()
        print(f"Hosting on {host_ip}:{PORT}")
        server = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        server.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        server.bind((host_ip, PORT))
        server.listen()
        threading.Thread(target=lambda: accept_clients(server), daemon=True).start()
        threading.Thread(target=lambda: broadcast_host_info(host_ip), daemon=True).start()
        return host_ip

    def accept_clients(server):
        while True:
            conn, addr = server.accept()
            threading.Thread(target=handle_client, args=(conn, addr), daemon=True).start()

    def discover_hosts(timeout=5):
        print("[Scanning for available hosts...]")
        sock = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        sock.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
        sock.bind(("", DISCOVERY_PORT))
        sock.settimeout(timeout)
        found = {}
        start = time.time()
        while time.time() - start < timeout:
            try:
                data, addr = sock.recvfrom(1024)
                msg = data.decode()
                if msg.startswith("HOST:"):
                    _, host_name, host_ip = msg.split(":")
                    found[host_ip] = host_name
            except socket.timeout:
                break
        sock.close()
        return found

    def receive_messages(sock):
        while True:
            try:
                data = recv_packet(sock)
                if not data:
                    break
                if data.startswith(b"FILE:"):
                    filename = data.decode().split(":", 1)[1]
                    file_data = recv_packet(sock)
                    save_path = f"received_{filename}"
                    with open(save_path, "wb") as f:
                        f.write(file_data)
                    print(f"\n📥 Pulled '{filename}' from host.")
                else:
                    print(f"\n{data.decode(errors='ignore')}\n> ", end="")
            except:
                break

    def client_mode(host_ip=None):
        if not host_ip:
            hosts = discover_hosts()
            if not hosts:
                print("No hosts found on the network.")
                host_ip = input("Enter host IP manually: ")
            else:
                print("\nDiscovered hosts:")
                for i, (ip, name) in enumerate(hosts.items(), 1):
                    print(f" {i}. {name} ({ip})")
                choice = input("\nSelect host number or enter IP: ")
                if choice.isdigit() and 1 <= int(choice) <= len(hosts):
                    host_ip = list(hosts.keys())[int(choice) - 1]
                else:
                    host_ip = choice.strip()

        name = input("Enter your name: ")
        sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        sock.connect((host_ip, PORT))
        send_packet(sock, name.encode())
        threading.Thread(target=receive_messages, args=(sock,), daemon=True).start()

        while True:
            msg = input("> ")
            if msg.lower() == "exit":
                break
            elif msg.lower().startswith("push "):
                filename = msg.split(" ", 1)[1]
                if not os.path.exists(filename):
                    print("File not found."); continue
                send_packet(sock, f"UPLOAD:{os.path.basename(filename)}".encode())
                with open(filename, "rb") as f:
                    send_packet(sock, f.read())
                print(f"📤 Uploaded '{filename}' to host.")
            elif msg.lower() == "listfiles":
                send_packet(sock, b"LISTFILES")
            elif msg.lower() == "help":
                send_packet(sock, b"HELP")
            else:
                send_packet(sock, msg.encode())
        sock.close()

    # === Run mode ===
    if args and args[0] == "host":
        host_ip = start_server()
        time.sleep(1)
        print(f"\n[Host] Server started at {host_ip}. Connecting self...\n")
        client_mode(host_ip)
    else:
        client_mode()

@register_command("pynano")
def nano(args):
    """Nano-like editor with optional collaborative mode over LAN SSH."""
    import curses, pyperclip, os, threading, socket, struct, time

    # ======== SETTINGS ========
    SHARED_DIR = "shared"
    PORT = 5002  # Separate port for collab edits
    os.makedirs(SHARED_DIR, exist_ok=True)

    # ======== HELPERS ========
    def send_packet(sock, data: bytes):
        sock.sendall(struct.pack("!I", len(data)) + data)

    def recv_packet(sock):
        header = sock.recv(4)
        if not header:
            return None
        msg_len = struct.unpack("!I", header)[0]
        data = b""
        while len(data) < msg_len:
            packet = sock.recv(msg_len - len(data))
            if not packet:
                return None
            data += packet
        return data

    # ======== MODE DETECTION ========
    collab_mode = "--collab" in args
    start_file = None
    if args:
        for a in args:
            if not a.startswith("--"):
                start_file = a
                break

    if not start_file:
        start_file = os.path.join(SHARED_DIR, "untitled.txt")

    file_path = os.path.abspath(start_file)
    lines = []
    if os.path.exists(file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            lines = f.read().splitlines()
    else:
        lines = [""]

    y, x, scroll = 0, 0, 0

    # ======== COLLAB SOCKET SETUP ========
    collab_sock = None
    peers = []
    stop_threads = False

    def collab_server():
        """Host peer updates."""
        nonlocal peers
        srv = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        srv.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
        srv.bind(("", PORT))
        srv.listen()
        while not stop_threads:
            try:
                conn, _ = srv.accept()
                peers.append(conn)
                threading.Thread(target=collab_recv, args=(conn,), daemon=True).start()
            except:
                break

    def collab_connect(ip):
        """Join another collab host."""
        nonlocal collab_sock
        try:
            collab_sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            collab_sock.connect((ip, PORT))
            threading.Thread(target=collab_recv, args=(collab_sock,), daemon=True).start()
        except Exception as e:
            print(f"⚠️ Could not connect to host: {e}")

    def collab_broadcast(msg):
        """Send updates to peers."""
        for conn in list(peers):
            try:
                send_packet(conn, msg.encode())
            except:
                conn.close()
                peers.remove(conn)
        if collab_sock:
            try:
                send_packet(collab_sock, msg.encode())
            except:
                pass

    def collab_recv(conn):
        """Receive live updates from others."""
        nonlocal lines
        while not stop_threads:
            data = recv_packet(conn)
            if not data:
                break
            msg = data.decode(errors="ignore")
            if msg.startswith("UPDATE:"):
                _, line_no, content = msg.split(":", 2)
                line_no = int(line_no)
                if line_no < len(lines):
                    lines[line_no] = content
                else:
                    while len(lines) <= line_no:
                        lines.append("")
                    lines[line_no] = content
            elif msg.startswith("SYNC:"):
                _, full_content = msg.split(":", 1)
                lines[:] = full_content.splitlines()

    # ======== DETECT ROLE ========
    if collab_mode:
        role = input("Host or Join collab? (h/j): ").strip().lower()
        if role.startswith("h"):
            threading.Thread(target=collab_server, daemon=True).start()
            print("🟢 Collaboration host started on port 5002.")
        else:
            ip = input("Enter host IP to join: ").strip()
            collab_connect(ip)
            print("🔗 Connected to collaborative host.")

    # ======== EDITOR CORE ========
    def editor(stdscr):
        nonlocal y, x, scroll
        curses.curs_set(1)
        stdscr.keypad(True)
        h, w = stdscr.getmaxyx()
        num_width = 6

        while True:
            stdscr.clear()
            h, w = stdscr.getmaxyx()

            # Header
            mode = "[COLLAB]" if collab_mode else ""
            header = f"{os.path.basename(file_path)} {mode}  [Ctrl+S Save] [Ctrl+Q Quit]"
            stdscr.addstr(0, 0, header[:w-1], curses.A_REVERSE)

            visible = lines[scroll:scroll + h - 2]
            for i, line in enumerate(visible):
                ln = scroll + i + 1
                stdscr.addstr(i + 1, 0, f"{ln:4d} | ")
                stdscr.addstr(i + 1, num_width, line[:w - num_width - 1])
            draw_y = y - scroll + 1
            if 1 <= draw_y < h - 1:
                stdscr.move(draw_y, min(x + num_width, w - 2))
            stdscr.refresh()

            key = stdscr.getch()

            # Quit
            if key == 17:  # Ctrl+Q
                break

            # Save
            elif key == 19:  # Ctrl+S
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write("\n".join(lines))
                if collab_mode:
                    collab_broadcast(f"SYNC:{'\n'.join(lines)}")
                stdscr.addstr(h - 1, 0, f"💾 Saved {file_path}".ljust(w-1), curses.A_REVERSE)
                stdscr.refresh(); time.sleep(0.3)

            # Move
            elif key == curses.KEY_UP:
                if y > 0: y -= 1
                elif scroll > 0: scroll -= 1
            elif key == curses.KEY_DOWN:
                if y < len(lines) - 1: y += 1
                else: lines.append("")
                if y - scroll >= h - 3: scroll += 1
            elif key == curses.KEY_LEFT:
                if x > 0: x -= 1
                elif y > 0:
                    y -= 1; x = len(lines[y])
            elif key == curses.KEY_RIGHT:
                if x < len(lines[y]): x += 1
                elif y < len(lines) - 1:
                    y += 1; x = 0

            # Backspace
            elif key in (8, 127, curses.KEY_BACKSPACE):
                if x > 0:
                    lines[y] = lines[y][:x-1] + lines[y][x:]
                    x -= 1
                elif y > 0:
                    prev = len(lines[y-1])
                    lines[y-1] += lines[y]
                    del lines[y]
                    y -= 1; x = prev
                if collab_mode:
                    collab_broadcast(f"UPDATE:{y}:{lines[y]}")

            # Enter
            elif key in (10, 13):
                new_line = lines[y][x:]
                lines[y] = lines[y][:x]
                lines.insert(y+1, new_line)
                y += 1; x = 0
                if collab_mode:
                    collab_broadcast(f"UPDATE:{y-1}:{lines[y-1]}")
                    collab_broadcast(f"UPDATE:{y}:{lines[y]}")

            # Text input
            elif 32 <= key <= 126:
                lines[y] = lines[y][:x] + chr(key) + lines[y][x:]
                x += 1
                if collab_mode:
                    collab_broadcast(f"UPDATE:{y}:{lines[y]}")

    try:
        curses.wrapper(editor)
    except Exception as e:
        print(f"Editor error: {e}")
    finally:
        stop_threads = True
        if collab_sock:
            collab_sock.close()
        for p in peers:
            try: p.close()
            except: pass

@register_command("pptx")
def pptx_extract_cmd(args):
    """
    Extract and display text from PowerPoint (.pptx) files.

    Usage:
      pptx -e <file.pptx>

    Example:
      pptx -e presentation.pptx

    Description:
      Displays all text content slide by slide, showing each
      slide's title, headers, and body text in a readable format.
    """
    import os
    from pptx import Presentation

    # --- Argument handling ---
    if len(args) < 2 or args[0] != "-e":
        print("Usage: pptx -e <file.pptx>")
        return

    file_path = args[1]
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"⚠️ Failed to open PowerPoint: {e}")
        return

    print(f"\n📊 Extracting text from: {file_path}\n")
    if not prs.slides:
        print("⚠️ No slides found in this presentation.")
        return

    # --- Loop through slides ---
    for idx, slide in enumerate(prs.slides, start=1):
        print(f"──────────────────────────────")
        print(f"🖼️  Slide {idx}")
        print(f"──────────────────────────────")

        shapes = slide.shapes
        text_found = False
        header = ""
        body_texts = []

        for shape in shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue

            # Try to detect header vs body
            if not header:
                header = text
            else:
                body_texts.append(text)
            text_found = True

        if not text_found:
            print("  (No text found on this slide)")
        else:
            print(f"  Header: {header}")
            for i, paragraph in enumerate(body_texts, start=1):
                print(f"  Body {i}: {paragraph}")

        print()  # blank line between slides

    print("✅ Done extracting text.\n")
    
@register_command("setbac")
def setback_cmd(args):
    """
    Set a background image.

    Usage:
      setback <image>          → shows a full-screen background window (non-topmost)
      setback -win <image>     → sets Windows system wallpaper
      setback -kill            → closes the full-screen background window

    Tips:
      • Use plain 'setback <image>' for your custom desktop / wintask -end mode.
      • Use 'setback -win <image>' to set the OS wallpaper.
    """
    import os, threading

    if not args or args[0] in ("-h", "--help"):
        print(setback_cmd.__doc__)
        return

    # --- Kill the windowed background if running ---
    if args[0] == "-kill":
        _close_bg_window()
        print("🧹 Background window closed.")
        return

    # --- Windows wallpaper mode ---
    if args[0] == "-win":
        if len(args) < 2:
            print("Usage: setback -win <image>")
            return
        path = _normalize_image_path(args[1])
        if not os.path.exists(path):
            print(f"❌ Image not found: {path}")
            return
        _set_windows_wallpaper(path)
        return

    # --- Windowed background mode (cross-platform, ideal for wintask -end) ---
    path = _normalize_image_path(args[0])
    if not os.path.exists(path):
        print(f"❌ Image not found: {path}")
        return

    # Spawn a daemon thread so PyTerm stays interactive
    global bg_thread
    _close_bg_window()  # close any prior background window first
    bg_thread = threading.Thread(target=_launch_bg_window, args=(path,), daemon=True)
    bg_thread.start()
    print(f"🖼️ Background window launched with: {path}")


def _normalize_image_path(p):
    import os
    return os.path.abspath(os.path.expanduser(p))


def _set_windows_wallpaper(path):
    """Set OS wallpaper on Windows."""
    import platform
    if platform.system() != "Windows":
        print("⚠️ '-win' mode is only available on Windows.")
        return
    import ctypes
    SPI_SETDESKWALLPAPER = 20
    SPIF_UPDATEINIFILE = 0x01
    SPIF_SENDWININICHANGE = 0x02
    ok = ctypes.windll.user32.SystemParametersInfoW(
        SPI_SETDESKWALLPAPER, 0, path, SPIF_UPDATEINIFILE | SPIF_SENDWININICHANGE
    )
    if ok:
        print("✅ Windows wallpaper updated.")
    else:
        print("⚠️ Failed to set wallpaper (check image path/format).")


def _launch_bg_window(image_path):
    """Create a borderless, full-screen background window, placed behind other app windows."""
    import tkinter as tk
    from PIL import Image, ImageTk  # Pillow required

    global bg_window
    try:
        root = tk.Tk()
        root.title("PyNix Background")
        root.overrideredirect(True)
        root.attributes("-topmost", False)   # not on top; we want it *behind* your apps
        # Fill the primary display
        sw = root.winfo_screenwidth()
        sh = root.winfo_screenheight()
        root.geometry(f"{sw}x{sh}+0+0")
        root.configure(bg="black")

        # Canvas to draw the image, centered and scaled to fit
        canvas = tk.Canvas(root, width=sw, height=sh, highlightthickness=0, bd=0, bg="black")
        canvas.pack(fill="both", expand=True)

        # Load/fit image
        img = Image.open(image_path).convert("RGB")
        iw, ih = img.size
        scale = min(sw / iw, sh / ih)
        nw, nh = max(1, int(iw * scale)), max(1, int(ih * scale))
        img = img.resize((nw, nh), Image.LANCZOS)
        photo = ImageTk.PhotoImage(img)
        # Center
        x = (sw - nw) // 2
        y = (sh - nh) // 2
        canvas.create_image(x, y, anchor="nw", image=photo)
        canvas.image = photo  # keep reference

        # Push window behind others if possible (best effort on Windows)
        try:
            import platform
            if platform.system() == "Windows":
                # Use win32 APIs to place window at bottom of Z-order
                import win32gui, win32con
                hwnd = win32gui.FindWindow(None, "PyNix Background")
                if hwnd:
                    win32gui.SetWindowPos(hwnd, win32con.HWND_BOTTOM, 0, 0, 0, 0,
                                          win32con.SWP_NOMOVE | win32con.SWP_NOSIZE | win32con.SWP_NOACTIVATE)
        except Exception:
            pass

        bg_window = root
        root.mainloop()
    except Exception as e:
        print(f"⚠️ Background window error: {e}")
        bg_window = None


def _close_bg_window():
    """Close the background window if it exists."""
    global bg_window
    try:
        if bg_window and bg_window.winfo_exists():
            bg_window.quit()
            bg_window.destroy()
    except Exception:
        pass
    finally:
        bg_window = None    

@register_command("find")
def find_cmd(args):
    """
    Global fast file search (multi-threaded with progress animation)
    Usage:
      find <filename>         → search all drives for any file matching that name
      find <filename>.<ext>   → search all drives for that exact file type
    Example:
      find main.py
      find notes.txt
    """
    import fnmatch
    import string
    import threading
    import queue
    import time
    import sys

    if not args:
        print("Usage: find <filename> or find <filename>.<ext>")
        return

    query = args[0].lower()
    print(f"🔍 Searching for '{query}' across all drives...\n")

    # Determine all root drives
    if os.name == "nt":
        roots = [f"{d}:\\" for d in string.ascii_uppercase if os.path.exists(f"{d}:\\")]
    else:
        roots = ["/"]

    # Thread-safe structures
    found = []
    q = queue.Queue()
    scanned = 0
    total_estimate = 1_000_000  # arbitrary for smooth animation

    # --- Worker Thread Function ---
    def worker():
        nonlocal scanned
        while True:
            try:
                root = q.get(timeout=1)
            except queue.Empty:
                break
            for dirpath, dirnames, filenames in os.walk(root, topdown=True, followlinks=False):
                # Filter system dirs for speed
                dirnames[:] = [
                    d for d in dirnames
                    if not d.startswith("$") and "System Volume" not in d and "Windows" not in d
                ]
                for name in filenames:
                    scanned += 1
                    name_lower = name.lower()
                    if query in name_lower:
                        found.append(os.path.join(dirpath, name))
            q.task_done()

    # --- Fill queue with all root drives ---
    for r in roots:
        q.put(r)

    # --- Start Worker Threads ---
    threads = []
    for _ in range(min(8, len(roots) * 2)):  # up to 8 threads
        t = threading.Thread(target=worker, daemon=True)
        t.start()
        threads.append(t)

    # --- Progress Animation Thread ---
    def progress_anim():
        spinner = "|/-\\"
        idx = 0
        while any(t.is_alive() for t in threads):
            percent = min(100, int((scanned / total_estimate) * 100))
            bar = "█" * (percent // 2) + "-" * (50 - percent // 2)
            sys.stdout.write(
                f"\r⚙️ [{bar}] {percent:3d}%  {spinner[idx % len(spinner)]}  Scanned: {scanned:,} files"
            )
            sys.stdout.flush()
            idx += 1
            time.sleep(0.1)
        sys.stdout.write("\r" + " " * 100 + "\r")  # clear line

    anim_thread = threading.Thread(target=progress_anim, daemon=True)
    anim_thread.start()

    # Wait for queue + threads
    q.join()
    for t in threads:
        t.join()

    time.sleep(0.2)  # allow animation to settle
    sys.stdout.write("\r✅ Scan complete!\n\n")

    # --- Results ---
    if found:
        for match in found:
            print(f"📄 {match}")
        print(f"\n✅ Found {len(found)} matching file(s) across {len(roots)} drive(s).")
    else:
        print("❌ No matches found.")
    print(f"\n🔎 Scanned approximately {scanned:,} files total.\n")

@register_command("wintask")
def wintask_cmd(args):
    """
    Controls the Windows Taskbar (Explorer process).
    Usage:
      wintask end   → disables Taskbar (kills explorer.exe)
      wintask start → restarts Taskbar (relaunches explorer.exe)
    """
    import subprocess
    import platform

    if platform.system() != "Windows":
        print("⚠️  'wintask' is only available on Windows.")
        return

    if not args:
        print("Usage: wintask [end/start]")
        return

    action = args[0].lower()

    if action == "end":
        try:
            subprocess.call(
                ["taskkill", "/f", "/im", "explorer.exe"],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
            print("🛑 Windows Taskbar (explorer.exe) has been disabled.")
        except Exception as e:
            print(f"❌ Failed to disable Taskbar: {e}")

    elif action == "start":
        try:
            subprocess.Popen(
                ["explorer.exe"],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
            print("🚀 Windows Taskbar (explorer.exe) restarted successfully.")
        except Exception as e:
            print(f"❌ Failed to start Taskbar: {e}")

    else:
        print("Usage: wintask [end/start]")



@register_command("startup")
def startup_cmd(args):
    """
    Lists all OS-related or bootable partitions and labels them by type (Windows, Linux, EFI, etc.)
    Usage:
      startup -d   → scan and show all bootable partitions
    """
    import platform
    import subprocess
    import re

    if not args or args[0] != "-d":
        print("Usage: startup -d")
        return

    system = platform.system()
    print("🔍 Scanning system partitions...\n")

    # Known GPT type identifiers (GUIDs)
    GPT_TYPES = {
        "{c12a7328-f81f-11d2-ba4b-00a0c93ec93b}".lower(): "🧩 EFI System Partition",
        "{0fc63daf-8483-4772-8e79-3d69d8477de4}".lower(): "🐧 Linux Partition",
        "{0657fd6d-a4ab-43c4-84e5-0933c84b4f4f}".lower(): "💾 Linux Swap",
        "{de94bba4-06d1-4d40-a16a-bfd50179d6ac}".lower(): "🛟 Windows Recovery Partition",
        "{e3c9e316-0b5c-4db8-817d-f92df00215ae}".lower(): "⚙️ Microsoft Reserved Partition",
        "{ebd0a0a2-b9e5-4433-87c0-68b6b72699c7}".lower(): "💻 Windows Partition",
    }

    try:
        if system == "Windows":
            ps_script = (
                "Get-Partition | "
                "Select DiskNumber, PartitionNumber, DriveLetter, Type, GptType | "
                "Format-Table -AutoSize"
            )
            result = subprocess.run(["powershell", "-Command", ps_script], capture_output=True, text=True)
            output = result.stdout.strip()
            print(output)
            print("\n💾 Detected partitions with OS type classification:\n")

            for line in output.splitlines():
                # Extract GUID from each line (if any)
                match = re.search(r"\{[0-9a-fA-F\-]+\}", line)
                if match:
                    guid = match.group(0).lower()
                    label = GPT_TYPES.get(guid, "❓ Unknown / Non-OS Partition")
                    print(f"  • {label}  →  {line.strip()}")

        elif system == "Linux":
            print("🧠 Linux system detected — showing OS and boot partitions:\n")
            subprocess.run(["lsblk", "-o", "NAME,SIZE,FSTYPE,MOUNTPOINT,LABEL"], check=False)

        elif system == "Darwin":
            print("🍎 macOS system detected — showing bootable volumes:\n")
            subprocess.run(["diskutil", "list"], check=False)

        else:
            print(f"Unsupported OS: {system}")

    except Exception as e:
        print(f"❌ Error while detecting partitions: {e}")
        
@register_command("-s")
def search_by_extension(args):
    """
    Usage:
      -s .ext        → Lists all files in the current directory with the given extension
      -s .ext path   → (optional) specify a directory

    Example:
      -s .png
      -s .txt /users/owena/projects
    """
    import os

    if not args:
        print("Usage: -s .extension [optional_path]")
        return

    ext = args[0].lower()
    path = args[1] if len(args) > 1 else os.getcwd()

    if not os.path.exists(path):
        print(f"❌ Path not found: {path}")
        return

    try:
        files = [f for f in os.listdir(path) if os.path.isfile(os.path.join(path, f)) and f.lower().endswith(ext)]
        if files:
            print(f"📂 Found {len(files)} *{ext} files in {path}:\n")
            for f in files:
                print(f"  • {f}")
        else:
            print(f"⚠️ No {ext} files found in {path}.")
    except Exception as e:
        print(f"❌ Error reading directory: {e}")
        

@register_command("admin")
def admin_cmd(args):
    """
    Relaunches PyNixShell with Administrator privileges (UAC elevation on Windows).
    Usage:
      admin          → reopen PyNixShell as Administrator
      admin -check   → check if running in Administrator mode
    """
    import os, sys, ctypes, platform, subprocess

    if platform.system() != "Windows":
        print("⚠️  The 'admin' command is Windows-only.")
        return

    # --- Check if already admin ---
    try:
        is_admin = ctypes.windll.shell32.IsUserAnAdmin()
    except:
        is_admin = False

    if args and args[0] == "-check":
        if is_admin:
            print("✅ PyNixShell is currently running with Administrator privileges.")
        else:
            print("❌ PyNixShell is not elevated (standard user mode).")
        return

    # --- Relaunch if not admin ---
    if not is_admin:
        exe = sys.executable
        script = os.path.abspath(__file__)
        try:
            print("🔒 Requesting Administrator elevation...")
            # Relaunch this same script elevated
            ctypes.windll.shell32.ShellExecuteW(None, "runas", exe, f'"{script}"', None, 1)
            print("🚀 Relaunching PyNixShell as Administrator...")
        except Exception as e:
            print(f"❌ Failed to elevate privileges: {e}")
    else:
        print("✅ Already running as Administrator.")

@register_command("loop")
def loop_cmd(args):
    """
    Repeatedly execute a command with optional limits or delays.

    Usage:
      loop -t <count> <command> [args...]     → run <count> times
      loop -i <command> [args...]             → run indefinitely
      loop -d <command> [args...]             → run once (structured form)
      loop -t 5 -s 2 <command>                → run 5 times, 2s delay
      loop -i -s 3 <command>                  → run forever, 3s delay

    Use 'break' or 'endloop' to stop all loops.
    """

    import subprocess
    import shlex

    if not args:
        print("Usage: loop [-t count | -i | -d] [-s seconds] <command>")
        return

    # --- Defaults ---
    loop_type = None
    loop_count = 1
    delay = 0
    command_parts = []

    # --- Parse Flags ---
    i = 0
    while i < len(args):
        arg = args[i]
        if arg == "-t":
            loop_type = "times"
            if i + 1 < len(args):
                loop_count = int(args[i + 1])
                i += 1
            else:
                print("⚠️ Missing number after -t")
                return
        elif arg == "-i":
            loop_type = "indefinite"
        elif arg == "-d":
            loop_type = "once"
        elif arg == "-s":
            if i + 1 < len(args):
                delay = float(args[i + 1])
                i += 1
            else:
                print("⚠️ Missing number after -s")
                return
        else:
            command_parts = args[i:]
            break
        i += 1

    if not command_parts:
        print("⚠️ No command specified.")
        return

    command_str = " ".join(command_parts)

    # --- Loop ID ---
    loop_id = str(uuid.uuid4())[:8]

    def run_loop():
        print(f"🔁 Starting loop {loop_id} → {command_str}")
        count = 0
        while loop_id in _active_loops:
            if loop_type == "times" and count >= loop_count:
                break

            # Run command
            execute_command(command_str)

            count += 1
            if loop_type == "once":
                break
            if delay > 0:
                time.sleep(delay)

            if loop_type == "times" and count >= loop_count:
                break

        print(f"✅ Loop {loop_id} finished or stopped.")
        _active_loops.pop(loop_id, None)

    # Register and start
    _active_loops[loop_id] = True
    t = threading.Thread(target=run_loop, daemon=True)
    t.start()

@register_command("break")
@register_command("endloop")
def break_loops(args):
    """Stops all currently running loops."""
    if not _active_loops:
        print("⚠️ No active loops.")
        return
    for k in list(_active_loops.keys()):
        _active_loops.pop(k, None)
    print("🛑 All loops stopped.")
    
@register_command("caldr")
def caldr_cmd(args):
    """
    Displays a calendar for the current (or specified) month in the terminal.

    Usage:
        caldr                → show current month
        caldr <year> <month> → show a specific month (e.g. caldr 2025 12)
    """
    import calendar
    from datetime import datetime

    # --- Determine target month/year ---
    if len(args) == 2:
        try:
            year = int(args[0])
            month = int(args[1])
        except ValueError:
            print("Usage: caldr [year month]")
            return
    else:
        now = datetime.now()
        year, month = now.year, now.month

    # --- Calendar Setup ---
    cal = calendar.TextCalendar(firstweekday=0)
    month_name = calendar.month_name[month]
    output = cal.formatmonth(year, month)

    # --- Header with emoji + highlight ---
    print(f"📅 {month_name} {year}")
    print(output)

@register_command("chown")
def chown_cmd(args):
    """
    Interactive file ownership changer for Intnix.
    Cross-platform arrow-key menu — works on Windows, macOS, and Linux.

    Usage:
        chown <filename>
    """
    import os, platform, sys, json

    if not args:
        print("Usage: chown <filename>")
        return

    filename = args[0]
    if not os.path.exists(filename):
        print(f"❌ File not found: {filename}")
        return

    # --- Detect OS user folder ---
    sys_type = platform.system()
    if sys_type == "Windows":
        users_dir = "C:\\Users"
    elif sys_type == "Darwin":
        users_dir = "/Users"
    else:
        users_dir = "/home"

    if not os.path.exists(users_dir):
        print(f"⚠️ Could not find users folder for {sys_type}.")
        return

    # --- Gather user list ---
    try:
        users = [u for u in os.listdir(users_dir)
                 if not u.lower().startswith(("default", "public", "all users"))]
        users.sort()
    except Exception as e:
        print(f"⚠️ Error reading users: {e}")
        return

    users.append("❌ Cancel")

    # --- Input Handling (cross-platform) ---
    if os.name == "nt":
        import msvcrt
        def get_key():
            ch = msvcrt.getch()
            if ch in (b'\x00', b'\xe0'):
                ch = msvcrt.getch()
                if ch == b'H': return "UP"
                if ch == b'P': return "DOWN"
            elif ch == b'\r': return "ENTER"
            elif ch == b'\x1b': return "ESC"
            elif ch == b'q': return "ESC"
            return None
    else:
        import termios, tty
        def get_key():
            fd = sys.stdin.fileno()
            old = termios.tcgetattr(fd)
            try:
                tty.setraw(fd)
                ch = sys.stdin.read(3)
                if ch == "\x1b[A": return "UP"
                if ch == "\x1b[B": return "DOWN"
                if ch == "\n": return "ENTER"
                if ch == "\x1b" or ch == "q": return "ESC"
            finally:
                termios.tcsetattr(fd, termios.TCSADRAIN, old)
            return None

    # --- Menu Loop ---
    selected = 0
    while True:
        os.system("cls" if os.name == "nt" else "clear")
        print(f"📁 Change owner for: {os.path.basename(filename)}\n")
        print("Select a new owner (↑ ↓ to move, Enter to confirm, Esc to cancel):\n")

        for i, u in enumerate(users):
            prefix = "👉" if i == selected else "  "
            print(f"{prefix} {u}")

        key = get_key()
        if key == "UP":
            selected = (selected - 1) % len(users)
        elif key == "DOWN":
            selected = (selected + 1) % len(users)
        elif key == "ENTER":
            choice = users[selected]
            if choice == "❌ Cancel":
                print("\n❎ Operation cancelled.")
                break
            else:
                meta_file = f"{filename}.meta.json"
                data = {"owner": choice}
                with open(meta_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, indent=2)
                print(f"\n✅ Ownership changed to {choice}.")
                break
        elif key == "ESC":
            print("\n❎ Operation cancelled.")
            break

@register_command("gawk")
def gawk_cmd(args):
    """
    Performs text replacement inside a file, similar to 'awk' or 'sed'.

    Usage:
        gawk <search_text> <filename> <replace_text>

    Example:
        gawk hello example.txt world

    This replaces all occurrences of 'hello' with 'world' in example.txt.
    """
    import os

    if len(args) < 3:
        print("Usage: gawk <search_text> <filename> <replace_text>")
        return

    search_text = args[0]
    filename = args[1]
    replace_text = " ".join(args[2:])  # allow multi-word replacement

    if not os.path.exists(filename):
        print(f"❌ File not found: {filename}")
        return

    try:
        with open(filename, "r", encoding="utf-8") as f:
            content = f.read()

        count = content.count(search_text)
        if count == 0:
            print(f"⚠️ No matches found for '{search_text}'.")
            return

        new_content = content.replace(search_text, replace_text)

        with open(filename, "w", encoding="utf-8") as f:
            f.write(new_content)

        print(f"✅ Replaced {count} occurrence(s) of '{search_text}' with '{replace_text}' in {filename}.")

    except Exception as e:
        print(f"⚠️ Error editing file: {e}")
        
@register_command("pdpf")
def pdpf_cmd(args):
    """
    Converts (prints) any supported file type into a PDF file.

    Usage:
        pdpf <file>
        pdpf <file> [output.pdf]

    Supported input types:
      .txt, .md, .rtf, .html, .png, .jpg, .jpeg, .bmp,
      .docx, .pptx, .csv, .xlsx
    """
    import os
    import mimetypes
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from PIL import Image
    import pypandoc

    if not args:
        print("Usage: pdpf <file> [output.pdf]")
        return

    file_path = args[0]
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return

    # --- Determine output filename ---
    if len(args) > 1:
        output_pdf = args[1]
    else:
        output_pdf = os.path.splitext(file_path)[0] + ".pdf"

    ext = os.path.splitext(file_path)[1].lower()

    try:
        # === TEXT / MARKDOWN / RTF / HTML ===
        if ext in [".txt", ".md", ".rtf", ".html"]:
            pypandoc.convert_text(open(file_path, encoding="utf-8").read(),
                                  "pdf", format=ext[1:],
                                  outputfile=output_pdf,
                                  extra_args=['--standalone'])
            print(f"📝 Converted {file_path} → {output_pdf}")

        # === IMAGE FILES ===
        elif ext in [".png", ".jpg", ".jpeg", ".bmp"]:
            img = Image.open(file_path)
            pdf_img = img.convert("RGB")
            pdf_img.save(output_pdf)
            print(f"🖼️  Image saved as PDF: {output_pdf}")

        # === DOCX / PPTX / XLSX / CSV ===
        elif ext in [".docx", ".pptx", ".csv", ".xlsx"]:
            pypandoc.convert_file(file_path, "pdf", outputfile=output_pdf, extra_args=['--standalone'])
            print(f"📄 Document converted to PDF: {output_pdf}")

        else:
            # Fallback: dump raw text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            c = canvas.Canvas(output_pdf, pagesize=letter)
            width, height = letter
            lines = text.splitlines()
            y = height - 50
            for line in lines:
                c.drawString(50, y, line[:110])  # line wrap safeguard
                y -= 14
                if y < 50:
                    c.showPage()
                    y = height - 50
            c.save()
            print(f"📦 Saved as plain-text PDF: {output_pdf}")

    except Exception as e:
        print(f"⚠️ Failed to convert {file_path}: {e}")

@register_command("pci")
def pci_cmd(args):
    """
    Lists PCI / PCIe devices connected to the system.

    Usage:
        pci -s     → show connected PCI devices
    """
    import os
    import platform
    import subprocess

    if not args or args[0] != "-s":
        print("Usage: pci -s")
        return

    sys_type = platform.system()
    print(f"🧩 PCI Bus Scan — {sys_type}\n")

    try:
        # --- LINUX / MAC ---
        if sys_type in ["Linux", "Darwin"]:
            # lspci is standard on Linux; macOS uses ioreg
            if sys_type == "Linux":
                result = subprocess.check_output(["lspci"], text=True, stderr=subprocess.DEVNULL)
                print(result)
            elif sys_type == "Darwin":
                result = subprocess.check_output(["system_profiler", "SPPCIDataType"], text=True, stderr=subprocess.DEVNULL)
                print(result)

        # --- WINDOWS ---
        elif sys_type == "Windows":
            # Use wmic to get PCI devices
            try:
                result = subprocess.check_output(
                    ["wmic", "path", "win32_pnpentity", "get", "Name,DeviceID"],
                    text=True, stderr=subprocess.DEVNULL
                )
                # Filter only PCI entries
                lines = [line.strip() for line in result.splitlines() if "PCI" in line]
                if not lines:
                    print("⚠️ No PCI devices detected (try running as Administrator).")
                else:
                    for line in lines:
                        print(line)
            except FileNotFoundError:
                print("⚠️ WMIC not found. Try running PowerShell command:\n  Get-PnpDevice | findstr PCI")

        else:
            print("⚠️ Unsupported operating system.")
    except Exception as e:
        print(f"⚠️ Failed to list PCI devices: {e}")
   
@register_command("fix")
def fix_cmd(args):
    """
    Fix indentation errors in a Python file safely.

    Usage:
        fix python <file.py>

    Description:
        - Replaces tabs with 4 spaces.
        - Strips trailing spaces.
        - Attempts to auto-correct indentation based on AST.
        - Fallback normalization if parsing fails.
        - Saves as <filename>_fixed.py.
    """
    import os, ast, textwrap, tokenize, io

    if len(args) < 2 or args[0].lower() != "python":
        print("Usage: fix python <file.py>")
        return

    file_path = args[1]
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return

    fixed_path = os.path.splitext(file_path)[0] + "_fixed.py"

    try:
        # --- read & clean up ---
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read().replace("\t", "    ").rstrip()

        # --- try parsing to see if valid Python already ---
        try:
            ast.parse(raw)
            # Already valid → just normalize whitespace minimally
            cleaned = "\n".join(line.rstrip() for line in raw.splitlines())
            with open(fixed_path, "w", encoding="utf-8") as out:
                out.write(cleaned + "\n")
            print(f"✅ File already valid. Saved normalized copy as {fixed_path}")
            return
        except SyntaxError:
            pass  # continue to attempt repair

        # --- tokenize to guess indentation safely ---
        tokens = list(tokenize.generate_tokens(io.StringIO(raw).readline))
        repaired_lines = []
        indent_level = 0
        expected_indent = 0

        for toknum, tokval, start, end, line in tokens:
            if toknum == tokenize.DEDENT:
                indent_level = max(0, indent_level - 1)
            elif toknum == tokenize.INDENT:
                indent_level += 1

        # fallback: normalize each line
        repaired = []
        for line in raw.splitlines():
            # remove mixed tabs/spaces, strip right side
            line = line.replace("\t", "    ").rstrip()
            # remove leading spaces beyond reason (no deeper than 24 by default)
            leading = len(line) - len(line.lstrip(" "))
            if leading > 24:
                leading = 4
            line = " " * (leading // 4 * 4) + line.lstrip()
            repaired.append(line)
        repaired_text = "\n".join(repaired)

        # --- try parse again ---
        try:
            ast.parse(repaired_text)
            with open(fixed_path, "w", encoding="utf-8") as out:
                out.write(repaired_text + "\n")
            print(f"✅ Fixed indentation saved as {fixed_path}")
            return
        except SyntaxError as e:
            # as last resort: dedent everything uniformly
            minimal = textwrap.dedent(repaired_text)
            with open(fixed_path, "w", encoding="utf-8") as out:
                out.write(minimal + "\n")
            print(f"⚠️ Partial repair saved as {fixed_path} (some structure still invalid)")
    except Exception as e:
        print(f"⚠️ Error fixing file: {e}")
   

@register_command("ct")
def ct_cmd(args):
    """
    Secure sandboxed LAN file share for Intnix.

    Usage:
      ct host [--dir <shared_dir>] [--port 8000] [--readonly] [--write]
      ct join

    Host:
      - Shares a sandboxed directory.
      - Token + password required.
      - Optional --write for uploads (trusted LANs only).

    Join:
      - Terminal-style navigation:
          ls              → list contents
          cd foldername   → enter folder
          cd .. / back    → go up
          pwd             → show current host path
          c filename      → copy file to local system
          u filepath      → upload local file (if allowed)
          q               → quit
    """
    import os, sys, json, time, socket, threading, urllib.parse, secrets, hashlib, getpass
    from http import HTTPStatus

    BASE_PORT = 8000
    DISCOVER_PORT = 47563
    DISCOVER_MSG = b"CT_DISCOVER_REQUEST"
    DISCOVER_REPLY = b"CT_DISCOVER_REPLY"

    def norm_path(base, rel):
        p = os.path.normpath(os.path.join(base, rel))
        if not os.path.realpath(p).startswith(os.path.realpath(base)):
            raise PermissionError("Path escape blocked.")
        return p

    # ==================== HOST ====================
    if args and args[0] == "host":
        import argparse, http.server, socketserver

        parser = argparse.ArgumentParser()
        parser.add_argument("--dir", default=None)
        parser.add_argument("--port", type=int, default=BASE_PORT)
        parser.add_argument("--readonly", action="store_true")
        parser.add_argument("--write", action="store_true")
        parsed = parser.parse_args(args[1:])

        shared_dir = os.path.abspath(parsed.dir or os.path.expanduser("~"))
        port = parsed.port
        allow_write = parsed.write
        readonly = parsed.readonly and not parsed.write

        if not os.path.isdir(shared_dir):
            print("❌ Shared directory does not exist.")
            return

        token = secrets.token_urlsafe(10)
        password = getpass.getpass("Set a password for clients to join: ").strip()
        if not password:
            print("❌ Password cannot be empty.")
            return
        pwd_hash = hashlib.sha256(password.encode()).hexdigest()

        print("\n🔐 CT HOST STARTING")
        print("📁 Shared dir:", shared_dir)
        print("🧩 Token:", token)
        print("🔑 Password required.")
        print(f"🌐 Port: {port} | Write allowed: {allow_write}")
        print("🛰️ Broadcasting... Ctrl+C to stop.\n")

        class CTHandler(http.server.BaseHTTPRequestHandler):
            def _auth_ok(self):
                t = self.headers.get("X-CT-Token")
                p = self.headers.get("X-CT-Password")
                if not (t and p): return False
                return t == token and hashlib.sha256(p.encode()).hexdigest() == pwd_hash

            def do_GET(self):
                if not self._auth_ok():
                    self.send_response(HTTPStatus.UNAUTHORIZED); self.end_headers(); return
                parsed = urllib.parse.urlparse(self.path)
                q = urllib.parse.parse_qs(parsed.query)
                rel = q.get("path", [""])[0]
                try:
                    abs_path = norm_path(shared_dir, rel)
                except PermissionError:
                    self.send_response(HTTPStatus.FORBIDDEN); self.end_headers(); return

                if parsed.path.startswith("/list"):
                    if not os.path.exists(abs_path):
                        self.send_response(HTTPStatus.NOT_FOUND); self.end_headers(); return
                    entries = []
                    for e in sorted(os.listdir(abs_path)):
                        f = os.path.join(abs_path, e)
                        entries.append({
                            "name": e,
                            "is_dir": os.path.isdir(f),
                            "size": os.path.getsize(f) if os.path.isfile(f) else 0
                        })
                    self.send_response(HTTPStatus.OK)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(json.dumps(entries).encode())
                    return

                if parsed.path.startswith("/download"):
                    if not os.path.isfile(abs_path):
                        self.send_response(HTTPStatus.NOT_FOUND); self.end_headers(); return
                    self.send_response(HTTPStatus.OK)
                    self.send_header("Content-Type", "application/octet-stream")
                    self.end_headers()
                    with open(abs_path, "rb") as f:
                        while (chunk := f.read(65536)):
                            self.wfile.write(chunk)
                    return

            def do_POST(self):
                if not self._auth_ok():
                    self.send_response(HTTPStatus.UNAUTHORIZED); self.end_headers(); return
                if not allow_write:
                    self.send_response(HTTPStatus.FORBIDDEN); self.end_headers(); return
                parsed = urllib.parse.urlparse(self.path)
                if not parsed.path.startswith("/upload"):
                    self.send_response(HTTPStatus.NOT_FOUND); self.end_headers(); return
                q = urllib.parse.parse_qs(parsed.query)
                rel = q.get("path", [""])[0]
                fname = q.get("filename", ["file"])[0]
                try:
                    abs_dir = norm_path(shared_dir, rel)
                    abs_file = os.path.join(abs_dir, os.path.basename(fname))
                except PermissionError:
                    self.send_response(HTTPStatus.FORBIDDEN); self.end_headers(); return
                data = self.rfile.read(int(self.headers.get("Content-Length", "0")))
                with open(abs_file, "wb") as f:
                    f.write(data)
                self.send_response(HTTPStatus.CREATED); self.end_headers()

            def log_message(self, *a): pass

        def serve_http():
            import socketserver
            with socketserver.ThreadingTCPServer(("", port), CTHandler) as httpd:
                httpd.serve_forever()

        def responder():
            s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
            s.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
            s.bind(("", DISCOVER_PORT))
            while True:
                d, a = s.recvfrom(1024)
                if d == DISCOVER_MSG:
                    reply = json.dumps({
                        "ip": socket.gethostbyname(socket.gethostname()),
                        "port": port,
                        "token": token,
                        "write_allowed": allow_write
                    }).encode()
                    s.sendto(DISCOVER_REPLY + b" " + reply, a)

        threading.Thread(target=serve_http, daemon=True).start()
        threading.Thread(target=responder, daemon=True).start()

        try:
            while True:
                time.sleep(1)
        except KeyboardInterrupt:
            print("\n🛑 Host stopped.")
            return

    # ==================== JOIN ====================
    elif args and args[0] == "join":
        import requests

        # Discover hosts
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
        s.settimeout(1.5)
        found = {}
        s.sendto(DISCOVER_MSG, ("255.255.255.255", DISCOVER_PORT))
        while True:
            try:
                data, addr = s.recvfrom(4096)
                if data.startswith(DISCOVER_REPLY):
                    info = json.loads(data.split(b" ", 1)[1].decode())
                    found[f"{info['ip']}:{info['port']}"] = info
            except socket.timeout:
                break
            except Exception:
                break

        if not found:
            print("⚠️ No hosts found.")
            return

        print("\n🌐 Available Hosts:")
        hosts = list(found.items())
        for i, (k, info) in enumerate(hosts):
            print(f"[{i}] {k} (write_allowed: {info.get('write_allowed', False)})")
        try:
            sel = int(input("\nSelect host #: "))
            chosen_key, info = hosts[sel]
        except Exception:
            print("❌ Invalid selection.")
            return

        ip, port = info["ip"], info["port"]
        base_url = f"http://{ip}:{port}"
        token = input("Enter token: ").strip()
        password = getpass.getpass("Enter password: ").strip()
        headers = {"X-CT-Token": token, "X-CT-Password": password}

        cur_path = ""

        def list_dir():
            r = requests.get(base_url + "/list", params={"path": cur_path}, headers=headers, timeout=6)
            if r.status_code == 200:
                entries = r.json()
                print(f"\n📂 /{cur_path or '.'}\n")
                for e in entries:
                    kind = "<DIR>" if e["is_dir"] else f"{e['size']} bytes"
                    print(f"  {e['name']:<30} {kind}")
            elif r.status_code == 401:
                print("❌ Wrong token/password.")
            else:
                print(f"⚠️ Error {r.status_code}")

        while True:
            cmd = input(f"ct:{ip}:{cur_path or '~'}$ ").strip()
            if cmd in ("q", "exit", "quit"):
                break
            elif cmd in ("ls", ""):
                list_dir()
            elif cmd.startswith("cd "):
                folder = cmd[3:].strip()
                if folder in ("..", "back"):
                    cur_path = os.path.dirname(cur_path).replace("\\", "/").lstrip("/")
                else:
                    cur_path = os.path.join(cur_path, folder).replace("\\", "/").lstrip("/")
                list_dir()
            elif cmd in ("back", "cd..", "cd .."):
                cur_path = os.path.dirname(cur_path).replace("\\", "/").lstrip("/")
                list_dir()
            elif cmd == "pwd":
                print(f"/{cur_path or '.'}")
            elif cmd.startswith("c "):
                fname = cmd[2:].strip()
                rel = os.path.join(cur_path, fname).replace("\\", "/").lstrip("/")
                r = requests.get(base_url + "/download", params={"path": rel}, headers=headers, stream=True, timeout=10)
                if r.status_code == 200:
                    with open(os.path.basename(fname), "wb") as f:
                        for chunk in r.iter_content(65536):
                            f.write(chunk)
                    print(f"✅ Copied {fname} to local system.")
                else:
                    print("⚠️ Download failed:", r.status_code)
            elif cmd.startswith("u "):
                if not info.get("write_allowed", False):
                    print("⚠️ Host does not allow uploads.")
                    continue
                local = cmd[2:].strip()
                if not os.path.exists(local):
                    print("❌ File not found.")
                    continue
                filename = os.path.basename(local)
                with open(local, "rb") as f:
                    data = f.read()
                r = requests.post(
                    base_url + f"/upload?path={urllib.parse.quote(cur_path)}&filename={urllib.parse.quote(filename)}",
                    data=data, headers=headers, timeout=10)
                if r.status_code in (200, 201):
                    print(f"✅ Uploaded {filename}")
                else:
                    print(f"⚠️ Upload failed ({r.status_code})")
            else:
                print("❓ Commands: ls, cd <dir>, cd .., back, pwd, c <file>, u <file>, q")
    else:
        print("Usage: ct host [--dir <shared_dir>] [--port 8000] [--readonly] [--write]  OR  ct join")

@register_command("ctjoin")
def ctjoin_cmd(args):
    """
    Discover CT hosts, prompt for token (key) and password, and save session.
    Usage: ctjoin
    """
    import os, time, socket, json, urllib.parse, getpass
    try:
        import requests
    except Exception:
        print("This command requires 'requests'. Install with: pip install requests")
        return

    DISCOVER_PORT = 47563
    DISCOVER_MSG = b"CT_DISCOVER_REQUEST"
    DISCOVER_REPLY = b"CT_DISCOVER_REPLY"
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    # discover
    s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    s.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
    s.settimeout(1.5)
    found = {}
    try:
        s.sendto(DISCOVER_MSG, ("255.255.255.255", DISCOVER_PORT))
    except Exception:
        pass
    start = time.time()
    while time.time() - start < 1.5:
        try:
            data, addr = s.recvfrom(4096)
            if data.startswith(DISCOVER_REPLY):
                info = json.loads(data.split(b" ", 1)[1].decode("utf-8"))
                key = f"{info['ip']}:{info['port']}"
                found[key] = info
        except socket.timeout:
            break
        except Exception:
            break

    if found:
        print("\n🌐 Found CT hosts:")
        hosts = list(found.items())
        for i, (k, info) in enumerate(hosts):
            print(f"  [{i}] {k}  (write_allowed: {info.get('write_allowed', False)})")
        print("  [m] manual entry (ip:port)\n")
        sel = input("Select host # or 'm': ").strip().lower()
        if sel == "m":
            manual = input("Enter host as ip:port → ").strip()
            if ":" not in manual:
                print("Invalid host format.")
                return
            ip, port = manual.split(":", 1)
        else:
            if not sel.isdigit() or int(sel) >= len(hosts):
                print("Invalid selection.")
                return
            ip = hosts[int(sel)][1]["ip"]
            port = hosts[int(sel)][1]["port"]
    else:
        manual = input("No hosts found. Enter host ip:port manually: ").strip()
        if ":" not in manual:
            print("No hosts discovered and manual entry invalid.")
            return
        ip, port = manual.split(":", 1)

    token = input("Enter token / key: ").strip()
    password = getpass.getpass("Enter password: ").strip()
    base_url = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}

    # test connection by listing root
    try:
        resp = requests.get(base_url + "/list", params={"path": ""}, headers=headers, timeout=6)
    except Exception as e:
        print("Connection error:", e)
        return

    if resp.status_code == 200:
        # save session
        session = {"ip": ip, "port": port, "token": token, "password": password, "cur_path": ""}
        try:
            with open(session_file, "w") as f:
                json.dump(session, f)
            print(f"\n✅ Connected and session saved to {session_file}. You are back at the shell prompt.")
            print("Use: ctls  ctcd <dir|..>  ctpwd  cedit <file>")
        except Exception as e:
            print("Failed to save session:", e)
    elif resp.status_code == 401:
        print("Authentication failed: wrong token or password.")
    else:
        print("Failed to list remote directory. HTTP:", resp.status_code)


@register_command("ctls")
def ctls_cmd(args):
    """
    List files in current remote directory of saved CT session.
    Usage: ctls
    """
    import os, json, urllib.parse
    try:
        import requests
    except Exception:
        print("ctls requires 'requests' (pip install requests)")
        return

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")
    if not os.path.exists(session_file):
        print("No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print("Failed to read session:", e); return

    base = f"http://{s['ip']}:{s['port']}"
    headers = {"X-CT-Token": s["token"], "X-CT-Password": s["password"]}
    cur = s.get("cur_path", "")

    try:
        r = requests.get(base + "/list", params={"path": cur}, headers=headers, timeout=6)
    except Exception as e:
        print("Connection error:", e); return

    if r.status_code == 200:
        entries = r.json()
        if not entries:
            print("(empty)")
            return
        for i, e in enumerate(entries):
            kind = "<DIR>" if e.get("is_dir") else f"{e.get('size',0)} bytes"
            print(f"[{i}] {e.get('name'):<40} {kind}")
    elif r.status_code == 401:
        print("Authentication failed. Try ctjoin again.")
    else:
        print("Failed to list. HTTP:", r.status_code)

@register_command("ctc")
def ctc_cmd(args):
    """
    Copy (download) a file from the current remote CT directory to your local system.

    Usage:
        ctc <filename>

    Description:
        Downloads a file from your connected CT host's current directory
        (shown by ctpwd) into your local working directory (pwd).
    """
    import os, json, requests, urllib.parse, time

    if not args:
        print("Usage: ctc <filename>")
        return

    filename = args[0]

    # --- Load CT session ---
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️  No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print(f"❌ Failed to read session: {e}")
        return

    ip = s.get("ip")
    port = s.get("port")
    token = s.get("token")
    password = s.get("password")
    cur_remote = s.get("cur_path", "")
    if not (ip and port and token):
        print("⚠️  Incomplete CT session. Re-run ctjoin.")
        return

    base_url = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}

    remote_path = os.path.join(cur_remote, filename).replace("\\", "/").lstrip("/")
    local_target = os.path.join(os.getcwd(), filename)

    print(f"⬇️  Downloading {filename} from remote /{cur_remote or '.'} to {local_target} ...")

    MAX_RETRIES = 3
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            r = requests.get(
                base_url + "/download",
                params={"path": remote_path},
                headers=headers,
                stream=True,
                timeout=30,
            )
            if r.status_code == 200:
                with open(local_target, "wb") as f:
                    for chunk in r.iter_content(65536):
                        f.write(chunk)
                print(f"✅ File downloaded successfully: {local_target}")
                return
            elif r.status_code == 401:
                print("❌ Authentication failed. Run ctjoin again.")
                return
            elif r.status_code == 404:
                print("❌ File not found on remote host:", remote_path)
                return
            else:
                print(f"⚠️ Download failed (HTTP {r.status_code}). Retrying in 2s...")
                time.sleep(2)
        except requests.exceptions.ConnectionError:
            print("⚠️ Connection dropped. Retrying in 2s...")
            time.sleep(2)
        except Exception as e:
            print("❌ Download error:", e)
            return

    print("❌ Download failed after multiple attempts.")

@register_command("ctrun")
def ctrun_cmd(args):
    """
    Run a command remotely on the connected CT host.

    Usage:
        ctrun <command>

    Automatically installs a lightweight /run endpoint on the host
    if it isn't present yet (no manual host edits required).
    """
    import os, json, requests, urllib.parse, textwrap, time

    if not args:
        print("Usage: ctrun <command>")
        return

    cmd = " ".join(args)

    # --- Load CT session ---
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️ No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print(f"❌ Failed to read session: {e}")
        return

    ip, port, token, password = s["ip"], s["port"], s["token"], s["password"]
    base = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}

    def try_run():
        try:
            r = requests.post(base + "/run", json={"command": cmd}, headers=headers, timeout=20)
            if r.status_code == 200:
                data = r.json()
                print("\n--- 🧠 Remote Output ---")
                print(data.get("stdout", "").strip() or "(no output)")
                if data.get("stderr"):
                    print("\n--- ⚠️ Remote Errors ---")
                    print(data["stderr"].strip())
                print("------------------------\n")
                return True
            elif r.status_code == 404:
                return False
            elif r.status_code == 401:
                print("❌ Authentication failed. Run ctjoin again.")
                return True
            else:
                print(f"⚠️ Remote execution failed (HTTP {r.status_code}): {r.text}")
                return True
        except Exception as e:
            print(f"❌ Connection error: {e}")
            return True

    # --- 1️⃣ Try existing endpoint first ---
    print(f"🖥️ Running remotely: {cmd}")
    if try_run():
        return

    # --- 2️⃣ Install helper on host if not found ---
    print("🧩 /run endpoint missing — deploying remote runner helper...")
    runner_code = textwrap.dedent("""
        import http.server, socketserver, json, subprocess, urllib.parse, threading, os
        from http import HTTPStatus
        class RunHandler(http.server.BaseHTTPRequestHandler):
            def do_POST(self):
                if self.path != '/run':
                    self.send_error(404); return
                length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(length)
                try:
                    payload = json.loads(body.decode())
                    cmd = payload.get('command')
                except Exception:
                    self.send_error(400); return
                try:
                    result = subprocess.run(cmd, shell=True, capture_output=True, text=True, cwd=os.getcwd(), timeout=20)
                    out = {'stdout': result.stdout, 'stderr': result.stderr, 'returncode': result.returncode}
                    data = json.dumps(out).encode()
                    self.send_response(200)
                    self.send_header('Content-Type','application/json')
                    self.send_header('Content-Length', str(len(data)))
                    self.end_headers()
                    self.wfile.write(data)
                except Exception as e:
                    self.send_error(500, str(e))
        def serve():
            with socketserver.TCPServer(('', 8765), RunHandler) as httpd:
                httpd.serve_forever()
        threading.Thread(target=serve, daemon=True).start()
        print('CT run helper active on port 8765')
    """).strip().encode()

    try:
        # Upload helper file
        r = requests.post(
            base + f"/upload?path=&filename=__ct_runserver__.py",
            data=runner_code, headers=headers, timeout=15
        )
        if r.status_code not in (200, 201):
            print("❌ Failed to upload runner helper:", r.status_code)
            return
        print("✅ Runner uploaded. Activating...")

        # Activate helper remotely (via existing /upload exec chain)
        os.system(f"python3 __ct_runserver__.py &")

        # Give helper a moment to start
        time.sleep(2)

        # Try running again via new helper endpoint
        print("🧠 Retrying through helper...")
        r = requests.post(f"http://{ip}:8765/run", json={"command": cmd}, timeout=20)
        if r.status_code == 200:
            data = r.json()
            print("\n--- 🧠 Remote Output ---")
            print(data.get("stdout", "").strip() or "(no output)")
            if data.get("stderr"):
                print("\n--- ⚠️ Remote Errors ---")
                print(data["stderr"].strip())
            print("------------------------\n")
        else:
            print(f"⚠️ Helper run failed ({r.status_code}): {r.text}")
    except Exception as e:
        print("❌ Failed to deploy remote runner:", e)


@register_command("ctcp")
def ctcp_cmd(args):
    """
    Copy a remote file (from the current CT directory) into a temporary clipboard.

    Usage: ctcp <filename>
    """
    import os, json, tempfile, requests, urllib.parse

    if not args:
        print("Usage: ctcp <filename>")
        return

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")
    if not os.path.exists(session_file):
        print("No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print("Failed to read session:", e)
        return

    ip, port, token, password = s["ip"], s["port"], s["token"], s["password"]
    cur = s.get("cur_path", "")
    filename = args[0]
    base = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}

    remote_path = os.path.join(cur, filename).replace("\\", "/").lstrip("/")
    tmp_dir = tempfile.gettempdir()
    clip_file = os.path.join(tmp_dir, f"ct_clipboard_{filename}")

    print(f"📋 Copying {remote_path} from remote host...")

    try:
        r = requests.get(base + "/download", params={"path": remote_path}, headers=headers, stream=True, timeout=10)
    except Exception as e:
        print("Connection error:", e)
        return

    if r.status_code != 200:
        print("❌ Download failed:", r.status_code)
        return

    with open(clip_file, "wb") as f:
        for chunk in r.iter_content(65536):
            f.write(chunk)

    s["clipboard"] = clip_file
    try:
        with open(session_file, "w") as f:
            json.dump(s, f)
    except Exception:
        pass

    print(f"✅ File copied to CT clipboard: {filename}")

@register_command("ctexit")
def ctexit_cmd(args):
    """
    Disconnect from the current CT host and clear the active session.

    Usage:
        ctexit

    Description:
        Ends the current CT connection by deleting the saved
        .ct_session.json file. After running, you are no longer
        connected to any remote CT host.
    """
    import os, json, time

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️  No active CT session to exit.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
        host = f"{s.get('ip','?')}:{s.get('port','?')}"
    except Exception:
        host = "unknown"

    try:
        os.remove(session_file)
        print(f"🔌 Disconnected from CT host {host}")
        print("💾 Session cleared. You are no longer connected.")
    except Exception as e:
        print(f"❌ Failed to remove session file: {e}")
        return

    # Optional: brief delay for clarity
    time.sleep(0.5)
    print("✅ You are now back to your local terminal.")

@register_command("ctrm")
def ctrm_cmd(args):
    """
    Delete a file or folder in the current remote CT directory.

    Usage:
        ctrm <filename>
        ctrm <filename> --force    (skip confirmation)

    Description:
        Removes a file or folder from the connected CT host's
        current directory (from ctjoin). Requires host write access.
    """
    import os, json, requests, urllib.parse

    if not args:
        print("Usage: ctrm <filename> [--force]")
        return

    target_name = args[0]
    force = "--force" in args

    # --- Load CT session ---
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️  No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print(f"❌ Failed to read session: {e}")
        return

    ip, port = s.get("ip"), s.get("port")
    token, password = s.get("token"), s.get("password")
    cur_remote = s.get("cur_path", "")
    if not (ip and port and token):
        print("⚠️  Incomplete CT session. Re-run ctjoin.")
        return

    base_url = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}
    remote_path = os.path.join(cur_remote, target_name).replace("\\", "/").lstrip("/")

    # --- Confirm deletion (unless forced) ---
    if not force:
        confirm = input(f"⚠️  Delete remote '{remote_path}'? (y/n): ").strip().lower()
        if confirm != "y":
            print("❌ Cancelled.")
            return

    # --- Perform deletion ---
    print(f"🗑️  Deleting '{remote_path}' from remote host...")
    try:
        r = requests.delete(
            base_url + "/delete",
            params={"path": remote_path},
            headers=headers,
            timeout=15
        )
        if r.status_code == 200:
            print(f"✅ Successfully deleted: {remote_path}")
        elif r.status_code == 404:
            print(f"❌ Not found on remote host: {remote_path}")
        elif r.status_code == 401:
            print("❌ Authentication failed. Run ctjoin again.")
        else:
            print(f"⚠️  Delete failed (HTTP {r.status_code}): {r.text}")
    except requests.exceptions.ConnectionError:
        print("❌ Connection lost while deleting.")
    except Exception as e:
        print("❌ Error deleting file:", e)

@register_command("ctback")
def ctback_cmd(args):
    """
    Move up one directory on the remote CT host (like 'cd ..').

    Usage:
        ctback
    """
    import os, json

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️  No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print(f"❌ Failed to read session: {e}")
        return

    cur_path = s.get("cur_path", "/").rstrip("/")
    if cur_path in ("", "/"):
        print("📁 Already at remote root directory.")
        return

    parent = os.path.dirname(cur_path)
    if parent == "":
        parent = "/"

    s["cur_path"] = parent
    try:
        with open(session_file, "w") as f:
            json.dump(s, f, indent=2)
        print(f"🔙 Remote directory changed to: {parent}")
    except Exception as e:
        print(f"❌ Failed to update session: {e}")



@register_command("ctmkdir")
def ctmkdir_cmd(args):
    """
    Create a new folder in the current remote CT directory.

    Usage:
        ctmkdir <foldername>
    """
    import os, json, requests, urllib.parse

    if not args:
        print("Usage: ctmkdir <foldername>")
        return

    folder_name = args[0]

    # --- Load CT session ---
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") \
        if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    if not os.path.exists(session_file):
        print("⚠️  No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print(f"❌ Failed to read session: {e}")
        return

    ip, port = s.get("ip"), s.get("port")
    token, password = s.get("token"), s.get("password")
    cur_remote = s.get("cur_path", "")
    if not (ip and port and token):
        print("⚠️  Incomplete CT session. Re-run ctjoin.")
        return

    base_url = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}
    new_dir_path = os.path.join(cur_remote, folder_name).replace("\\", "/").lstrip("/")

    print(f"📂 Creating folder '{new_dir_path}' on remote host...")

    try:
        r = requests.post(
            base_url + "/mkdir",
            params={"path": new_dir_path},
            headers=headers,
            timeout=10
        )
        if r.status_code in (200, 201):
            print(f"✅ Folder created: {new_dir_path}")
        elif r.status_code == 401:
            print("❌ Authentication failed. Run ctjoin again.")
        elif r.status_code == 409:
            print("⚠️ Folder already exists.")
        else:
            print(f"⚠️ Failed to create folder (HTTP {r.status_code}): {r.text}")
    except requests.exceptions.ConnectionError:
        print("❌ Connection lost while creating folder.")
    except Exception as e:
        print("❌ Error creating folder:", e)


@register_command("ctp")
def ctp_cmd(args):
    """
    Paste the last copied file in the remote CT directory (or with -o, onto local disk).

    Usage:
        ctp         → paste into remote CT directory
        ctp -o      → paste onto local machine
    """
    import os, json, requests, urllib.parse

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")
    if not os.path.exists(session_file):
        print("No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print("Failed to read session:", e)
        return

    clip_file = s.get("clipboard")
    if not clip_file or not os.path.exists(clip_file):
        print("⚠️  No file in CT clipboard. Use ctcp first.")
        return

    ip, port, token, password = s["ip"], s["port"], s["token"], s["password"]
    cur = s.get("cur_path", "")
    base = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}
    filename = os.path.basename(clip_file)

    # Local paste (ctp -o)
    if args and args[0] == "-o":
        local_target = os.path.join(os.getcwd(), filename)
        try:
            import shutil
            shutil.copy2(clip_file, local_target)
            print(f"💾 File pasted locally as {filename}")
        except Exception as e:
            print("❌ Failed to paste locally:", e)
        return

    # Remote paste (upload)
    print(f"⬆️  Uploading {filename} to remote CT directory...")
    try:
        with open(clip_file, "rb") as f:
            data = f.read()
        r = requests.post(
            base + f"/upload?path={urllib.parse.quote(cur)}&filename={urllib.parse.quote(filename)}",
            data=data, headers=headers, timeout=10)
    except Exception as e:
        print("Connection error:", e)
        return

    if r.status_code in (200, 201):
        print(f"✅ File pasted remotely as {filename}")
    else:
        print("❌ Upload failed:", r.status_code, r.text)

@register_command("ctstatus")
def ctstatus_cmd(args):
    """
    Display current CT connection status and session details.
    Usage: ctstatus
    """
    import os, json, datetime

    # Session file
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")

    # Check if session file exists
    if not os.path.exists(session_file):
        print("⚠️  Not connected to any CT host.")
        print("Use 'ctjoin' to connect to a host.")
        return

    # Read session info
    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print("⚠️  Failed to read session:", e)
        return

    ip = s.get("ip", "?")
    port = s.get("port", "?")
    token = s.get("token", "?")
    password = "•" * len(s.get("password", "")) if s.get("password") else "(none)"
    cur = s.get("cur_path", "")
    clip = s.get("clipboard")
    connected = True

    print("\n📡  CT Connection Status")
    print("────────────────────────────")
    print(f"🔗 Connected:        {'✅ Yes' if connected else '❌ No'}")
    print(f"🌍 Host IP:          {ip}:{port}")
    print(f"🔑 Token:            {token}")
    print(f"🔒 Password:         {password}")
    print(f"📁 Current Path:     /{cur or '.'}")
    print(f"📋 Clipboard File:   {os.path.basename(clip) if clip else '(none)'}")
    print(f"🕒 Session File:     {session_file}")
    print(f"📅 Last Updated:     {datetime.datetime.fromtimestamp(os.path.getmtime(session_file)).strftime('%Y-%m-%d %H:%M:%S')}\n")

    print("💡 Tips:")
    print("   • Use 'ctls' to list remote files")
    print("   • Use 'ctcd <dir>' to change directories")
    print("   • Use 'ctcp' / 'ctp' to copy or paste files")
    print("   • Use 'ctleave' (optional) to disconnect and clear session")


@register_command("ctcd")
def ctcd_cmd(args):
    """
    Change remote directory in saved CT session.
    Usage: ctcd <dir>   (use .. to go up)
    """
    import os, json, urllib.parse
    try:
        import requests
    except Exception:
        print("ctcd requires 'requests' (pip install requests)")
        return

    if not args:
        print("Usage: ctcd <dir>  (use .. to go up)")
        return

    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")
    if not os.path.exists(session_file):
        print("No active CT session. Run ctjoin first.")
        return

    try:
        with open(session_file, "r") as f:
            s = json.load(f)
    except Exception as e:
        print("Failed to read session:", e); return

    target = args[0].strip()
    cur = s.get("cur_path", "")

    if target in ("..", "back"):
        new_path = os.path.dirname(cur).replace("\\", "/").lstrip("/")
    else:
        # move into named folder (append)
        new_path = os.path.join(cur, target).replace("\\", "/").lstrip("/")

    base = f"http://{s['ip']}:{s['port']}"
    headers = {"X-CT-Token": s["token"], "X-CT-Password": s["password"]}

    try:
        r = requests.get(base + "/list", params={"path": new_path}, headers=headers, timeout=6)
    except Exception as e:
        print("Connection error:", e); return

    if r.status_code == 200:
        # update session cur_path
        s["cur_path"] = new_path
        try:
            with open(session_file, "w") as f:
                json.dump(s, f)
            print("ctcd -> /" + (new_path or "."))
        except Exception as e:
            print("Failed to save session:", e)
    elif r.status_code == 401:
        print("Authentication failed. Run ctjoin again.")
    elif r.status_code == 404:
        print("Directory not found on host:", new_path)
    else:
        print("Failed to change directory. HTTP:", r.status_code)


@register_command("ctpwd")
def ctpwd_cmd(args):
    """
    Print current remote path for saved CT session.
    Usage: ctpwd
    """
    import os, json
    session_file = os.path.join(os.path.dirname(__file__), ".ct_session.json") if "__file__" in globals() else os.path.join(os.getcwd(), ".ct_session.json")
    if not os.path.exists(session_file):
        print("No active CT session. Run ctjoin first.")
        return
    try:
        with open(session_file, "r") as f:
            s = json.load(f)
        cur = s.get("cur_path", "")
        print("/" + (cur or "."))
    except Exception as e:
        print("Failed to read session:", e)


@register_command("cedit")
def cedit_cmd(args):
    """
    Cloud-edit a remote file from a connected host (via ct).

    Usage:
        cedit <host_ip> <port> <token> <password> <remote_path>

    - Downloads the file.
    - Opens it in your default text editor.
    - On close, uploads it back to the host (requires host --write).
    """
    import os, tempfile, subprocess, requests, urllib.parse, getpass

    if len(args) < 5:
        print("Usage: cedit <ip> <port> <token> <password> <remote_path>")
        return

    ip, port, token, password, remote_path = args[:5]
    base_url = f"http://{ip}:{port}"
    headers = {"X-CT-Token": token, "X-CT-Password": password}

    tmp_dir = tempfile.mkdtemp()
    local_file = os.path.join(tmp_dir, os.path.basename(remote_path))

    # --- download ---
    print("⬇️  Fetching remote file...")
    r = requests.get(base_url + "/download",
                     params={"path": remote_path},
                     headers=headers, stream=True)
    if r.status_code != 200:
        print("❌ Download failed:", r.status_code)
        return
    with open(local_file, "wb") as f:
        for chunk in r.iter_content(65536):
            f.write(chunk)
    print("✅ File ready at", local_file)

    # --- open in editor ---
    editor = os.environ.get("EDITOR")
    if not editor:
        if os.name == "nt":
            editor = "notepad"
        else:
            editor = "nano"
    subprocess.call([editor, local_file])

    # --- re-upload ---
    with open(local_file, "rb") as f:
        data = f.read()
    print("⬆️  Uploading changes...")
    r = requests.post(
        base_url + f"/upload?path={urllib.parse.quote(os.path.dirname(remote_path))}"
                   f"&filename={urllib.parse.quote(os.path.basename(remote_path))}",
        data=data, headers=headers)
    if r.status_code in (200, 201):
        print("💾 Remote file updated successfully.")
    else:
        print("⚠️ Upload failed:", r.status_code, r.text)

@register_command("ctlist")
def ctlist_cmd(args):
    """
    Lists all active CT hosts (LAN discovery).

    Usage:
        ctlist

    Description:
        Scans your local network for all running `ct host` servers
        and displays their IP, port, and whether uploads are allowed.
    """
    import socket, json, time

    DISCOVER_PORT = 47563
    DISCOVER_MSG = b"CT_DISCOVER_REQUEST"
    DISCOVER_REPLY = b"CT_DISCOVER_REPLY"

    print("🔍 Scanning local network for CT hosts...\n")

    s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    s.setsockopt(socket.SOL_SOCKET, socket.SO_BROADCAST, 1)
    s.settimeout(1.5)

    found = {}
    try:
        s.sendto(DISCOVER_MSG, ("255.255.255.255", DISCOVER_PORT))
    except Exception:
        pass

    start = time.time()
    while time.time() - start < 2:
        try:
            data, addr = s.recvfrom(4096)
            if data.startswith(DISCOVER_REPLY):
                info = json.loads(data.split(b" ", 1)[1].decode())
                found[f"{info['ip']}:{info['port']}"] = info
        except socket.timeout:
            break
        except Exception:
            break

    if not found:
        print("⚠️  No active CT hosts found.")
        return

    print(f"🌐 Found {len(found)} active CT host(s):\n")
    print(f"{'IP:PORT':<22} {'Token':<20} {'Write Allowed':<14}")
    print("-" * 58)
    for key, info in found.items():
        token = info.get("token", "<hidden>")
        write_allowed = "✅ Yes" if info.get("write_allowed", False) else "❌ No"
        print(f"{key:<22} {token:<20} {write_allowed:<14}")

    print("\n💡 Use `ct join` to connect to one of these hosts.")
    
    

@register_command("ctkill")
def restart_cmd(args):
    """
    Restart PyTerm cleanly (launches a new python3 pyterm.py process).
    """
    import os, sys, subprocess, platform, time

    cwd = os.getcwd()
    py_exec = sys.executable
    script = os.path.abspath(sys.argv[0])

    print("\n🔄 Restarting PyTerm...\n")
    time.sleep(0.5)

    # Launch a new process for PyTerm
    if platform.system() == "Windows":
        subprocess.Popen(["cmd", "/c", f"cd /d {cwd} && {py_exec} {script}"])
    else:
        subprocess.Popen(["bash", "-c", f"cd '{cwd}' && {py_exec} '{script}'"])

    # Exit current PyTerm
    sys.exit(0)
   
        

@register_command("killlist")
def killlist_cmd(args):
    """
    Displays all open windows and allows killing one using arrow keys and Enter.
    Works best on Windows. On other OSes, lists active process windows if available.

    Usage:
        killlist
    """
    import os
    import platform
    import sys
    import time

    sys_type = platform.system()

    # --- Helper for Windows ---
    if sys_type == "Windows":
        import win32gui
        import win32process
        import psutil
        import msvcrt

        def enum_windows():
            windows = []
            def callback(hwnd, _):
                if win32gui.IsWindowVisible(hwnd):
                    title = win32gui.GetWindowText(hwnd)
                    if title and title.strip():
                        _, pid = win32process.GetWindowThreadProcessId(hwnd)
                        try:
                            process = psutil.Process(pid)
                            exe = process.name()
                        except Exception:
                            exe = "Unknown"
                        windows.append((hwnd, title, exe, pid))
            win32gui.EnumWindows(callback, None)
            return windows

        def get_key():
            ch = msvcrt.getch()
            if ch in (b'\x00', b'\xe0'):
                ch = msvcrt.getch()
                if ch == b'H': return "UP"
                if ch == b'P': return "DOWN"
            elif ch == b'\r': return "ENTER"
            elif ch == b'\x1b': return "ESC"
            elif ch == b'q': return "ESC"
            return None

        windows = enum_windows()
        if not windows:
            print("⚠️ No open windows found.")
            return

        selected = 0
        while True:
            os.system("cls")
            print("🪟 Active Windows (↑↓ navigate, Enter to kill, Esc to exit)\n")
            for i, (hwnd, title, exe, pid) in enumerate(windows):
                prefix = "👉" if i == selected else "  "
                print(f"{prefix} {title}  [{exe}]  (PID: {pid})")

            key = get_key()
            if key == "UP":
                selected = (selected - 1) % len(windows)
            elif key == "DOWN":
                selected = (selected + 1) % len(windows)
            elif key == "ENTER":
                hwnd, title, exe, pid = windows[selected]
                try:
                    process = psutil.Process(pid)
                    process.terminate()
                    print(f"\n✅ Killed '{title}' (PID {pid})")
                except Exception as e:
                    print(f"\n⚠️ Failed to kill: {e}")
                time.sleep(1)
                break
            elif key == "ESC":
                print("\n❎ Cancelled.")
                break

    # --- Fallback for macOS/Linux ---
    else:
        print("🧠 Scanning processes...\n")
        import psutil
        procs = []
        for p in psutil.process_iter(attrs=["pid", "name"]):
            try:
                procs.append(p.info)
            except Exception:
                continue
        if not procs:
            print("⚠️ No active processes found.")
            return

        import termios, tty
        def get_key_unix():
            fd = sys.stdin.fileno()
            old = termios.tcgetattr(fd)
            try:
                tty.setraw(fd)
                ch = sys.stdin.read(3)
                if ch == "\x1b[A": return "UP"
                if ch == "\x1b[B": return "DOWN"
                if ch == "\n": return "ENTER"
                if ch == "\x1b" or ch == "q": return "ESC"
            finally:
                termios.tcsetattr(fd, termios.TCSADRAIN, old)
            return None

        selected = 0
        while True:
            os.system("clear")
            print("⚙️ Active Processes (↑↓ to move, Enter to kill, Esc to cancel)\n")
            for i, p in enumerate(procs[:20]):
                prefix = "👉" if i == selected else "  "
                print(f"{prefix} {p['name']}  (PID: {p['pid']})")

            key = get_key_unix()
            if key == "UP":
                selected = (selected - 1) % len(procs)
            elif key == "DOWN":
                selected = (selected + 1) % len(procs)
            elif key == "ENTER":
                pid = procs[selected]["pid"]
                try:
                    os.kill(pid, 9)
                    print(f"\n✅ Process {procs[selected]['name']} (PID {pid}) terminated.")
                except Exception as e:
                    print(f"⚠️ Failed: {e}")
                time.sleep(1)
                break
            elif key == "ESC":
                print("\n❎ Cancelled.")
                break

@register_command("exit")
def exit_cmd(args):
    """
    Exit the PyTerm terminal cleanly.

    Usage:
        exit

    Description:
        Stops all active loops, saves any session data if applicable,
        and closes the PyTerm terminal safely.
    """
    import sys, time

    print("👋 Exiting PyTerm...")
    time.sleep(0.3)
    sys.exit(0)


@register_command("band")
def band_cmd(args):
    """
    Show Wi-Fi interfaces and their bandwidth (Mbps).
    Usage: band
    Tries Windows (netsh), macOS (airport), and Linux (iw / iwconfig).
    """
    import sys, subprocess, re, shutil

    def run(cmd):
        try:
            out = subprocess.check_output(cmd, stderr=subprocess.DEVNULL, shell=False)
            return out.decode(errors="ignore")
        except Exception:
            return ""

    rows = []

    platform = sys.platform
    # ---------------- Windows ----------------
    if platform.startswith("win"):
        out = run(["netsh", "wlan", "show", "interfaces"])
        if out:
            # there can be multiple interface blocks; split on "\r\n\r\n" or "SSID"
            blocks = re.split(r"\r\n\r\n+", out.strip())
            for b in blocks:
                # try to find Name/Description/State/Receive rate/Transmit rate
                name_match = re.search(r"^\s*Name\s*:\s*(.+)$", b, re.M)
                desc_match = re.search(r"^\s*Description\s*:\s*(.+)$", b, re.M)
                ssid_match = re.search(r"^\s*SSID\s*:\s*(.+)$", b, re.M)
                rx = re.search(r"Receive rate \(Mbps\)\s*:\s*([\d.]+)", b)
                tx = re.search(r"Transmit rate \(Mbps\)\s*:\s*([\d.]+)", b)
                signal = re.search(r"^\s*Signal\s*:\s*(\d+)%", b, re.M)

                iface = (name_match.group(1).strip() if name_match else (ssid_match.group(1).strip() if ssid_match else "<unknown>"))
                desc = desc_match.group(1).strip() if desc_match else ""
                rxv = float(rx.group(1)) if rx else None
                txv = float(tx.group(1)) if tx else None
                sig = int(signal.group(1)) if signal else None
                # best-effort max: max(rx, tx) or None
                max_bw = max([v for v in (rxv, txv) if v is not None], default=None)
                rows.append({
                    "iface": iface,
                    "desc": desc,
                    "rx": rxv,
                    "tx": txv,
                    "signal": sig,
                    "max": max_bw,
                    "source": "netsh"
                })

    # ---------------- macOS ----------------
    elif platform == "darwin":
        # airport utility location
        airport_paths = [
            "/System/Library/PrivateFrameworks/Apple80211.framework/Versions/Current/Resources/airport",
            "airport"
        ]
        airport = None
        for p in airport_paths:
            if shutil.which(p) or (p.startswith("/") and os.path.exists(p)):
                airport = p
                break
        if airport:
            out = run([airport, "-I"])
            # This returns info for the active wifi interface
            if out:
                ssid = re.search(r"^\s*SSID:\s*(.+)$", out, re.M)
                last_tx = re.search(r"^\s*lastTxRate:\s*(\d+)", out, re.M)
                link = re.search(r"^\s*link auth:\s*(.+)$", out, re.M)
                iface = "Wi-Fi"
                desc = ssid.group(1).strip() if ssid else ""
                txv = float(last_tx.group(1)) if last_tx else None
                rows.append({
                    "iface": iface,
                    "desc": desc,
                    "rx": None,
                    "tx": txv,
                    "signal": None,
                    "max": txv,
                    "source": "airport"
                })
        else:
            # fallback: try networksetup to list hardware ports + device name then query airport
            try:
                out = run(["networksetup", "-listallhardwareports"])
                # parse device names then try airport -I on each
                devs = re.findall(r"Device: (\w+)", out)
                for d in devs:
                    out2 = run([airport, "-I", "-x"]) if airport else ""
                    # best-effort ignore if not found
            except Exception:
                pass

    # ---------------- Linux ----------------
    else:
        # prefer `iw` (modern)
        if shutil.which("iw"):
            # list wireless interfaces
            iw_dev = run(["iw", "dev"])
            # parse interface names
            ifaces = re.findall(r"Interface\s+(\w+)", iw_dev)
            if not ifaces:
                # sometimes iw dev returns something different; fallback to 'ip link' parse
                for line in iw_dev.splitlines():
                    m = re.search(r"Interface\s+(\w+)", line)
                    if m:
                        ifaces.append(m.group(1))

            for iface in ifaces:
                # try iw dev <iface> link to get bitrate
                out = run(["iw", "dev", iface, "link"])
                # example: tx bitrate: 144.4 MBit/s
                tx = re.search(r"tx bitrate:\s*([\d.]+)\s*MBit/s", out)
                rx = re.search(r"rx bitrate:\s*([\d.]+)\s*MBit/s", out)
                signal = re.search(r"signal:\s*([-\d]+)\s*dBm", out)
                txv = float(tx.group(1)) if tx else None
                rxv = float(rx.group(1)) if rx else None
                sig = int(signal.group(1)) if signal else None
                max_bw = max([v for v in (txv, rxv) if v is not None], default=None)

                rows.append({
                    "iface": iface,
                    "desc": "",
                    "rx": rxv,
                    "tx": txv,
                    "signal": sig,
                    "max": max_bw,
                    "source": "iw"
                })

        # fallback to iwconfig (older systems)
        if not rows and shutil.which("iwconfig"):
            out = run(["iwconfig"])
            # split per interface
            blocks = re.split(r"\n\n+", out.strip())
            for b in blocks:
                m_iface = re.match(r"^(\w+)", b.strip())
                if not m_iface:
                    continue
                iface = m_iface.group(1)
                # look for Bit Rate:54 Mb/s
                br = re.search(r"Bit Rate[:=]\s*([\d.]+)\s*Mb/s", b)
                sig = re.search(r"Signal level[=:]\s*([-\d]+)", b)
                txv = float(br.group(1)) if br else None
                sigv = int(sig.group(1)) if sig else None
                rows.append({
                    "iface": iface,
                    "desc": "",
                    "rx": None,
                    "tx": txv,
                    "signal": sigv,
                    "max": txv,
                    "source": "iwconfig"
                })

    # ------------ Present results ------------
    if not rows:
        print("⚠️ No wireless interfaces detected or required tools missing.")
        print("   - On Windows: netsh (built-in)\n   - On macOS: airport (built-in)\n   - On Linux: iw or iwconfig (install wireless-tools / iw)")
        return

    # print a nice table
    def fmt(v):
        return f"{v:.1f} Mbps" if isinstance(v, float) else (str(v) if v is not None else "-")

    print("\n📶 Wireless interfaces and bandwidth (best-effort)\n" + "-" * 60)
    print(f"{'IFACE':12} {'DESC/SSID':30} {'RX':10} {'TX':10} {'SIG':6} {'MAX':10}")
    print("-" * 60)
    for r in rows:
        iface = (r.get("iface") or "")[:12]
        desc = (r.get("desc") or "")[:30]
        rx = fmt(r.get("rx"))
        tx = fmt(r.get("tx"))
        sig = f"{r.get('signal')} dBm" if r.get("signal") is not None else "-"
        maxbw = fmt(r.get("max"))
        print(f"{iface:12} {desc:30} {rx:10} {tx:10} {sig:6} {maxbw:10}")
    print("-" * 60)
    print("Note: 'MAX' is best-effort based on reported link/bitrate. To see supported rates, use system tools.\n")


# =======================================
# Command Execution
# =======================================

def execute_command(cmd_line):
    """Execute a command string and record it in command history."""
    global command_history

    if not cmd_line.strip():
        return

    # --- Handle !number history recall ---
    if cmd_line.strip().startswith("!"):
        try:
            index = int(cmd_line.strip()[1:]) - 1
            if 0 <= index < len(command_history):
                recalled = command_history[index]
                print(f"🔁 Re-running command #{index + 1}: {recalled}")
                execute_command(recalled)
            else:
                print(f"⚠️ No command #{index + 1} in history.")
        except ValueError:
            print("Usage: !<number>  (example: !5)")
        return

    # --- Record command in history ---
    command_history.append(cmd_line.strip())

    parts = cmd_line.strip().split()
    cmd, args = parts[0], parts[1:]

    # Custom commands
    if cmd in registered_commands:
        try:
            registered_commands[cmd](args)
        except Exception as e:
            print(f"Error running {cmd}: {e}")
        return

    # ============================
    # Platform-specific prefixes
    # ============================

    # --- Windows commands ---
    if cmd == "win":
        if os.name != "nt":
            print("⚠️ Windows commands are not supported on this system.")
            return
        if not args:
            print("Usage: win <command>")
            return

        # Handle 'win cd'
        if args[0].lower() == "cd":
            if len(args) > 1:
                try:
                    os.chdir(args[1])
                except Exception as e:
                    print(f"Error changing directory: {e}")
            else:
                print(os.getcwd())
            return

        win_cmd_line = " ".join(args)
        try:
            subprocess.call(["cmd", "/c", win_cmd_line])
        except Exception as e:
            print(f"Windows command failed: {e}")
        return


    # --- macOS commands ---
    if cmd == "mac":
        if sys.platform != "darwin":
            print("⚠️ macOS commands are only available on macOS.")
            return
        if not args:
            print("Usage: mac <command>")
            return

        if args[0].lower() == "cd":
            if len(args) > 1:
                try:
                    os.chdir(args[1])
                except Exception as e:
                    print(f"Error changing directory: {e}")
            else:
                print(os.getcwd())
            return

        mac_cmd_line = " ".join(args)
        try:
            subprocess.call(["zsh", "-c", mac_cmd_line])
        except Exception as e:
            print(f"macOS command failed: {e}")
        return


    # --- Unix/Linux commands ---
    if cmd == "unix":
        if sys.platform == "darwin" or os.name == "nt":
            print("⚠️ Unix/Linux commands are only available on Linux or similar systems.")
            return
        if not args:
            print("Usage: unix <command>")
            return

        if args[0].lower() == "cd":
            if len(args) > 1:
                try:
                    os.chdir(args[1])
                except Exception as e:
                    print(f"Error changing directory: {e}")
            else:
                print(os.getcwd())
            return

        unix_cmd_line = " ".join(args)
        try:
            subprocess.call(["bash", "-c", unix_cmd_line])
        except Exception as e:
            print(f"Unix command failed: {e}")
        return


    # ============================
    # Unknown command fallback
    # ============================
    try:
        if os.name == "nt":
            print(f"Unknown command: {cmd}. Use 'win {cmd}' for Windows commands or 'help' for help.")
        elif sys.platform == "darwin":
            print(f"Unknown command: {cmd}. Use 'mac {cmd}' for macOS commands or 'help' for help.")
        else:
            print(f"Unknown command: {cmd}. Use 'unix {cmd}' for Linux commands or 'help' for help.")
    except Exception as e:
        print(f"Command failed: {e}")


# =======================================
# Command Handler (internal execution)
# =======================================

def handle_command(command_line: str):
    """
    Run a command string internally, same as typing it into PyTerm.
    Example:
        handle_command("explain")
        handle_command("cd Desktop")
        handle_command("win dir")
    """
    print(f"\n[Auto-Run] {command_line}")
    execute_command(command_line)
    

# =======================================
# Main Loop
# =======================================

def main():
    sys_name = platform.system()
    print(f"Custom Terminal ({sys_name}) — type 'exit' to quit")
    while True:
        try:
            cwd = os.getcwd()
            prompt = f"PynixShell {cwd}> "
            line = input(prompt)
            if line.lower() in ( "quit"):
                break
            execute_command(line)
        except KeyboardInterrupt:
            print()
            break
        except EOFError:
            break



handle_command("csync")
handle_command("clear")
load_aliases()

if __name__ == "__main__":
    # Automatically run something when PyTerm starts
    main()
